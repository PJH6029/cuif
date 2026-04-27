from __future__ import annotations

import base64
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from cuif_eval.pptx.render import render_pptx_previews


TASK_DIR = Path(__file__).resolve().parent
ARTIFACTS = TASK_DIR / "artifacts"
INPUTS = ARTIFACTS / "inputs"
GOLD_DIR = ARTIFACTS / "gold"
MOCK = TASK_DIR / "mock_outputs"

NAVY = RGBColor(23, 32, 51)
BLUE = RGBColor(31, 78, 121)
STEEL = RGBColor(88, 103, 120)
PALE_BLUE = RGBColor(231, 241, 251)
MINT = RGBColor(232, 247, 240)
RED = RGBColor(193, 59, 54)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(249, 251, 253)
LIGHT_GRAY = RGBColor(225, 231, 238)


def ensure_dirs() -> None:
    for path in [ARTIFACTS, INPUTS, GOLD_DIR, MOCK / "turn1", MOCK / "final"]:
        path.mkdir(parents=True, exist_ok=True)


def prs_16x9() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    return prs


def add_box(slide, x, y, w, h, fill=WHITE, line=LIGHT_GRAY, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.25)
    return shape


def add_text(slide, text, x, y, w, h, *, size=18, color=NAVY, bold=False, fill=None, line=None, margin=0.08, name=None):
    if fill is None and line is None:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        shape = add_box(slide, x, y, w, h, fill=fill or WHITE, line=line or LIGHT_GRAY)
    if name:
        shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_metric(slide, text, x, y, *, fill=PALE_BLUE):
    shape = add_text(slide, text, x, y, 3.35, 0.76, size=20, bold=True, fill=fill, line=RGBColor(174, 193, 211), name=text)
    return shape


def add_background(slide) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = OFF_WHITE


def add_protected_intro(slide) -> None:
    add_background(slide)
    add_text(slide, "Nimbus Reliability Operations", 0.62, 0.46, 7.8, 0.55, size=30, bold=True, color=NAVY)
    add_text(slide, "Executive context for April service review", 0.64, 1.08, 7.0, 0.36, size=15, color=STEEL)
    add_text(
        slide,
        "Scope: support, SRE, and comms readiness\nPrimary KPI: restore severity-2 ownership discipline\nDo not edit: Nimbus context invariant",
        0.78,
        2.0,
        5.9,
        1.7,
        size=18,
        fill=WHITE,
        line=LIGHT_GRAY,
    )
    add_text(
        slide,
        "Protected slide for collateral-damage checks",
        7.35,
        2.08,
        4.9,
        0.58,
        size=18,
        bold=True,
        fill=MINT,
        line=RGBColor(111, 181, 160),
    )
    add_text(slide, "Revision target is slide 2 only.", 7.42, 2.88, 4.1, 0.42, size=16, color=STEEL)


def add_runbook_slide(slide) -> None:
    add_background(slide)
    add_text(slide, "Runbook Handoff", 0.62, 0.48, 4.8, 0.52, size=30, bold=True)
    add_text(slide, "Do not edit: runbook appendix invariant", 0.66, 1.1, 5.8, 0.4, size=15, color=STEEL)
    add_text(
        slide,
        "Escalation ladder\n1. Triage lead confirms ownership\n2. SRE opens mitigation channel\n3. Comms drafts customer update",
        0.72,
        2.0,
        5.7,
        2.2,
        size=18,
        fill=WHITE,
        line=LIGHT_GRAY,
    )
    add_text(
        slide,
        "Weekend coverage window\nFriday 18:00 to Monday 09:00\nOwner rotation remains unchanged.",
        7.05,
        2.0,
        4.9,
        1.6,
        size=18,
        fill=WHITE,
        line=LIGHT_GRAY,
    )


def add_audit_slide(slide) -> None:
    add_background(slide)
    add_text(slide, "Benchmark Audit Trail", 0.62, 0.48, 5.5, 0.52, size=30, bold=True)
    add_text(
        slide,
        "Do not edit: incident response benchmark invariant\nSeed slide 2 is the only editable target.\nPreserve source narrative and non-target slides.",
        0.74,
        1.72,
        7.6,
        1.56,
        size=18,
        fill=WHITE,
        line=LIGHT_GRAY,
    )
    add_text(slide, "CUIF check focus: targeted edit + preservation", 0.76, 4.18, 6.7, 0.5, size=18, color=STEEL)


def add_seed_target_slide(slide) -> None:
    add_background(slide)
    add_text(slide, "Reliability Operations Review", 0.58, 0.38, 7.0, 0.48, size=28, bold=True)
    add_text(slide, "Seed slide: current layout is crowded and needs the screenshot annotations applied.", 0.61, 0.92, 8.8, 0.35, size=13, color=STEEL)
    add_metric(slide, "MTTA 8 min", 0.82, 1.48)
    add_metric(slide, "Backlog 42", 8.62, 1.22)
    add_metric(slide, "Coverage 92%", 4.55, 5.82)
    add_text(
        slide,
        "Incident trend\nP1 3 | P2 11 | P3 28\nSmall chart panel: expand and make this the analytical focus.",
        0.82,
        3.55,
        5.25,
        1.88,
        size=17,
        fill=WHITE,
        line=RGBColor(167, 181, 198),
        name="incident_trend_panel_seed",
    )
    add_text(
        slide,
        "SEV-2 queue risk\nEscalation owner missing\nAnnotate hotfix owner by Friday",
        8.44,
        2.0,
        3.74,
        1.72,
        size=17,
        fill=RGBColor(255, 239, 239),
        line=RED,
        name="sev2_queue_risk_seed",
    )
    add_text(
        slide,
        "Unsorted handoff notes\n- incidents trending down\n- capacity gap weekends\n- execs need an owner view",
        6.35,
        4.32,
        5.1,
        1.68,
        size=15,
        fill=WHITE,
        line=LIGHT_GRAY,
    )


def add_turn1_target_slide(slide, *, final: bool = False) -> None:
    add_background(slide)
    title_size = 32 if final else 28
    title_color = BLUE if final else NAVY
    add_text(slide, "Reliability Operations Review", 0.58, 0.34, 7.2, 0.52, size=title_size, bold=True, color=title_color, name="slide2_title")
    add_text(slide, "Edited slide: annotated screenshot applied to the target dashboard.", 0.61, 0.92, 8.1, 0.34, size=13, color=STEEL)
    add_metric(slide, "MTTA 8 min", 0.68, 1.32)
    add_metric(slide, "Backlog 42", 4.92, 1.32)
    add_metric(slide, "Coverage 92%", 9.16, 1.32)
    chart = add_text(
        slide,
        "Incident trend\nP1 incidents: 3\nP2 incidents: 11\nP3 incidents: 28",
        0.72,
        2.44,
        7.55,
        3.1,
        size=19,
        fill=WHITE,
        line=RGBColor(116, 140, 166),
        name="incident_trend_panel",
    )
    chart.line.width = Pt(2.0)
    for idx, (label, width, color) in enumerate(
        [
            ("P1", 0.72, RGBColor(91, 141, 203)),
            ("P2", 1.72, RGBColor(68, 124, 185)),
            ("P3", 3.36, BLUE),
        ]
    ):
        y = 4.56 + idx * 0.28
        add_box(slide, 3.42, y, width, 0.15, fill=color, line=color, radius=False)
        add_text(slide, label, 2.92, y - 0.08, 0.35, 0.22, size=9, color=STEEL)
    if final:
        add_text(
            slide,
            "Caption: queue risk is concentrated in SEV-2 ownership gaps.",
            0.82,
            5.72,
            7.3,
            0.42,
            size=15,
            color=STEEL,
            name="queue_risk_caption",
        )
        add_text(
            slide,
            "Action owners\nSRE: Priya\nSupport: Marcus\nComms: Lina",
            8.86,
            2.54,
            3.58,
            1.82,
            size=17,
            fill=MINT,
            line=RGBColor(111, 181, 160),
            name="action_owners",
        )
    add_text(
        slide,
        "SEV-2 queue risk\nEscalation owner missing\nAnnotate hotfix owner by Friday",
        8.86,
        5.08,
        3.58,
        1.42,
        size=17,
        fill=RGBColor(255, 239, 239),
        line=RED,
        name="sev2_queue_risk",
    )


def build_deck(stage: str) -> Presentation:
    prs = prs_16x9()
    for _ in range(4):
        prs.slides.add_slide(prs.slide_layouts[6])
    add_protected_intro(prs.slides[0])
    if stage == "seed":
        add_seed_target_slide(prs.slides[1])
    elif stage == "turn1":
        add_turn1_target_slide(prs.slides[1], final=False)
    elif stage == "final":
        add_turn1_target_slide(prs.slides[1], final=True)
    else:
        raise ValueError(stage)
    add_runbook_slide(prs.slides[2])
    add_audit_slide(prs.slides[3])
    return prs


def write_revision_notes() -> None:
    (INPUTS / "revision_notes.txt").write_text(
        "\n".join(
            [
                "Scenario: executive reliability operations review.",
                "Only slide 2 should be edited. Slides 1, 3, and 4 contain benchmark/provenance text and must be preserved.",
                "The annotated screenshot marks three intended layout repairs: align the metric tiles across the top row, expand the Incident trend panel into the left-middle region, and move the SEV-2 queue risk callout to the lower-right.",
                "The final turn adds executive handoff details without moving the turn-1 chart or damaging protected slides.",
                "",
            ]
        ),
        encoding="utf-8",
    )


def write_annotated_screenshot() -> None:
    preview_dir = INPUTS / "_seed_preview"
    preview = render_pptx_previews(ARTIFACTS / "seed.pptx", preview_dir, "seed_for_annotation")
    images = preview.get("images", [])
    if len(images) < 2:
        raise RuntimeError(f"Could not render slide 2 for annotation: {preview}")
    slide2_png = INPUTS / "slide2_seed_screenshot.png"
    shutil.copy2(images[1], slide2_png)
    encoded = base64.b64encode(slide2_png.read_bytes()).decode("ascii")
    (INPUTS / "annotated_slide2_screenshot.svg").write_text(
        f"""<svg xmlns="http://www.w3.org/2000/svg" width="1152" height="648" viewBox="0 0 1920 1080">
  <defs>
    <marker id="arrow" markerWidth="14" markerHeight="14" refX="11" refY="5" orient="auto" markerUnits="strokeWidth">
      <path d="M0,0 L12,5 L0,10 z" fill="#c13b36"/>
    </marker>
    <style>
      .note {{ font-family: "Comic Sans MS", "Bradley Hand", "Segoe Print", Arial, sans-serif; fill: #c13b36; font-size: 38px; font-weight: 700; }}
      .thin {{ fill: none; stroke: #c13b36; stroke-width: 8; stroke-linecap: round; stroke-linejoin: round; }}
      .target {{ fill: rgba(193,59,54,.04); stroke: #c13b36; stroke-width: 7; stroke-dasharray: 20 16; rx: 18; }}
    </style>
  </defs>
  <image href="data:image/png;base64,{encoded}" x="0" y="0" width="1920" height="1080"/>
  <rect class="target" x="92" y="188" width="1738" height="128"/>
  <path class="thin" d="M258 152 C520 70 1044 76 1426 144" marker-end="url(#arrow)"/>
  <text class="note" x="170" y="92">Line up all three metric tiles across this top row</text>
  <rect class="target" x="92" y="350" width="1110" height="470"/>
  <path class="thin" d="M376 913 C402 780 516 636 684 520" marker-end="url(#arrow)"/>
  <text class="note" x="104" y="972">Make "Incident trend" much larger here</text>
  <rect class="target" x="1268" y="728" width="514" height="252"/>
  <path class="thin" d="M1702 704 C1658 644 1582 636 1504 694" marker-end="url(#arrow)"/>
  <text class="note" x="1016" y="656">Move risk callout to lower-right</text>
  <path class="thin" d="M1448 992 C1288 982 1168 982 1052 952" marker-end="url(#arrow)"/>
  <text class="note" x="990" y="1020">Keep all risk text; don't touch slides 1, 3, or 4</text>
</svg>
""",
        encoding="utf-8",
    )
    shutil.rmtree(preview_dir, ignore_errors=True)


def main() -> None:
    ensure_dirs()
    write_revision_notes()
    seed = build_deck("seed")
    seed.save(ARTIFACTS / "seed.pptx")
    turn1 = build_deck("turn1")
    turn1.save(GOLD_DIR / "turn1.pptx")
    final = build_deck("final")
    final.save(GOLD_DIR / "final.pptx")
    write_annotated_screenshot()
    shutil.copy2(GOLD_DIR / "turn1.pptx", MOCK / "turn1" / "result.pptx")
    shutil.copy2(GOLD_DIR / "final.pptx", MOCK / "final" / "result.pptx")


if __name__ == "__main__":
    main()
