from __future__ import annotations

import shutil
import subprocess
import tempfile
import urllib.request
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from cuif_eval.pptx.render import render_pptx_previews


TASK_DIR = Path(__file__).resolve().parent
ARTIFACTS = TASK_DIR / "artifacts"
INPUTS = ARTIFACTS / "inputs"
MOCK = TASK_DIR / "mock_outputs"

PDF_URL = "https://arxiv.org/pdf/1706.03762"
PDF_PATH = INPUTS / "attention_is_all_you_need.pdf"
FIGURE1_PATH = INPUTS / "figure1_transformer_architecture.png"
FIGURE2_PATH = INPUTS / "figure2_attention_mechanisms.png"
FORMULA_IMAGE_PATH = INPUTS / "attention_formula_crop.png"

NAVY = RGBColor(23, 32, 51)
BLUE = RGBColor(31, 78, 121)
INK = RGBColor(38, 42, 56)
STEEL = RGBColor(86, 102, 122)
MIST = RGBColor(237, 244, 249)
MINT = RGBColor(229, 247, 240)
GOLD_COLOR = RGBColor(233, 180, 71)
PAPER = RGBColor(250, 248, 242)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(218, 225, 232)
RED = RGBColor(198, 61, 55)
PURPLE = RGBColor(230, 224, 244)


def ensure_dirs() -> None:
    for path in [ARTIFACTS, INPUTS, ARTIFACTS / "gold", MOCK / "turn1", MOCK / "turn2", MOCK / "final"]:
        path.mkdir(parents=True, exist_ok=True)


def download_pdf() -> None:
    if PDF_PATH.exists() and PDF_PATH.stat().st_size > 500_000:
        return
    request = urllib.request.Request(PDF_URL, headers={"User-Agent": "cuif-poc-task-generator/0.1"})
    with urllib.request.urlopen(request, timeout=60) as response:
        PDF_PATH.write_bytes(response.read())


def render_pdf_page(page_number: int, out_dir: Path) -> Path:
    prefix = out_dir / f"page_{page_number}"
    subprocess.run(
        ["pdftoppm", "-png", "-r", "160", "-f", str(page_number), "-l", str(page_number), str(PDF_PATH), str(prefix)],
        check=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
    )
    matches = sorted(out_dir.glob(f"{prefix.name}-*.png"))
    if not matches:
        raise FileNotFoundError(f"pdftoppm did not create page image for page {page_number}")
    return matches[0]


def trim_white(image: Image.Image, *, pad: int = 16) -> Image.Image:
    gray = image.convert("L")
    mask = gray.point(lambda value: 255 if value < 248 else 0)
    bbox = mask.getbbox()
    if bbox is None:
        return image
    left, top, right, bottom = bbox
    return image.crop((max(0, left - pad), max(0, top - pad), min(image.width, right + pad), min(image.height, bottom + pad)))


def crop_pdf_assets() -> None:
    if FIGURE1_PATH.exists() and FIGURE2_PATH.exists() and FORMULA_IMAGE_PATH.exists():
        return
    with tempfile.TemporaryDirectory(prefix="cuif_attention_pdf_") as tmp:
        tmp_dir = Path(tmp)
        page3 = Image.open(render_pdf_page(3, tmp_dir)).convert("RGB")
        page4 = Image.open(render_pdf_page(4, tmp_dir)).convert("RGB")

        figure1 = trim_white(page3.crop((410, 135, 980, 965)))
        figure1.save(FIGURE1_PATH)

        figure2 = trim_white(page4.crop((300, 150, 1085, 580)), pad=18)
        figure2.save(FIGURE2_PATH)

        formula = trim_white(page4.crop((480, 1015, 1005, 1088)), pad=12)
        formula.save(FORMULA_IMAGE_PATH)


def prs_16x9() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_background(slide, *, fill=PAPER) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill


def add_box(slide, x, y, w, h, *, fill=WHITE, line=GRAY, radius=True, width=1.2, name=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(width)
    return shape


def add_text(slide, text, x, y, w, h, *, size=18, color=NAVY, bold=False, fill=None, line=None, margin=0.05, name=None):
    shape = add_box(slide, x, y, w, h, fill=fill, line=line or fill, name=name) if fill is not None else slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_multiline(slide, title, lines, x, y, w, h, *, title_color=BLUE, fill=WHITE, line=GRAY, name=None):
    shape = add_box(slide, x, y, w, h, fill=fill, line=line, width=1.1, name=name)
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.14)
    frame.margin_right = Inches(0.14)
    frame.margin_top = Inches(0.10)
    frame.margin_bottom = Inches(0.08)
    p = frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = title_color
    for line in lines:
        p = frame.add_paragraph()
        p.space_before = Pt(4)
        p.level = 0
        r = p.add_run()
        r.text = line
        r.font.size = Pt(12)
        r.font.color.rgb = INK
    return shape


def add_pill(slide, text, x, y, w, *, fill=MIST, line=GRAY, color=BLUE, name=None):
    return add_text(slide, text, x, y, w, 0.34, size=10.5, color=color, bold=True, fill=fill, line=line, margin=0.04, name=name)


def add_footer(slide, text="SNUPI Reading Group | arXiv:1706.03762") -> None:
    add_text(slide, text, 0.62, 7.02, 6.4, 0.22, size=8.5, color=STEEL)


def add_header(slide, title, eyebrow=None, *, final_style=False, name=None) -> None:
    add_box(slide, 0, 0, 13.333, 0.16, fill=BLUE, line=BLUE, radius=False, width=0)
    add_text(slide, title, 0.62, 0.40, 8.8, 0.45, size=32 if final_style else 28, color=BLUE if final_style else NAVY, bold=True, name=name)
    if eyebrow:
        add_pill(slide, eyebrow, 10.3, 0.44, 2.3, fill=WHITE, line=GRAY)


def add_picture_fit(slide, image_path: Path, x, y, w=None, h=None, *, name=None):
    kwargs = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    pic = slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), **kwargs)
    if name:
        pic.name = name
    return pic


def add_seed_slide1(slide) -> None:
    add_background(slide)
    add_header(slide, "SNUPI Paper Seminar Template", "SEED")
    add_text(slide, "Paper PDF: attention_is_all_you_need.pdf", 0.72, 1.34, 5.2, 0.4, size=18, bold=True, color=INK)
    add_multiline(
        slide,
        "Raw seminar notes",
        [
            "Turn this template into review slides for a reading group.",
            "Use the PDF, not web summaries, for the paper figure and attention formula.",
            "Keep the protected logistics slide unchanged.",
        ],
        0.72,
        2.05,
        5.7,
        2.15,
        fill=WHITE,
    )
    add_multiline(
        slide,
        "Expected deck shape",
        ["Cover", "Architecture figure", "Attention mechanism + formula", "Verdict and discussion", "Protected logistics"],
        7.05,
        2.05,
        4.95,
        2.15,
        fill=MIST,
    )
    add_footer(slide)


def add_seed_slide2(slide) -> None:
    add_background(slide)
    add_header(slide, "Architecture Figure Workspace", "ANNOTATE ME")
    add_text(slide, "Paste/crop Figure 1 from the paper PDF here.", 0.82, 1.18, 5.6, 0.34, size=14, color=STEEL)
    add_box(slide, 0.78, 1.65, 5.65, 4.55, fill=WHITE, line=BLUE, width=1.6, name="seed_figure_placeholder")
    add_text(slide, "FIGURE 1\nTransformer model architecture", 1.18, 3.35, 4.2, 0.70, size=22, color=BLUE, bold=True)
    add_multiline(
        slide,
        "Notes to convert",
        ["Explain encoder and decoder halves.", "Call out multi-head attention bridge.", "Caption must cite Figure 1."],
        7.05,
        1.72,
        5.0,
        2.15,
        fill=MIST,
    )
    add_multiline(
        slide,
        "Revision target",
        ["Follow annotated_slide2_layout.png on turn 2.", "Do not change other slides while fixing slide 2."],
        7.05,
        4.20,
        5.0,
        1.5,
        fill=WHITE,
    )
    add_footer(slide)


def add_seed_slide3(slide) -> None:
    add_background(slide)
    add_header(slide, "Attention Formula Workspace", "PDF EXTRACT")
    add_box(slide, 0.78, 1.45, 5.4, 3.2, fill=WHITE, line=GRAY, name="seed_figure2_placeholder")
    add_text(slide, "Figure 2 from PDF", 2.25, 2.68, 2.8, 0.36, size=20, color=STEEL, bold=True)
    add_box(slide, 6.85, 1.45, 5.55, 3.2, fill=MIST, line=GRAY, name="seed_formula_placeholder")
    add_text(slide, "Scaled dot-product formula from PDF", 7.25, 2.72, 4.45, 0.36, size=18, color=STEEL, bold=True)
    add_footer(slide)


def add_seed_slide4(slide) -> None:
    add_background(slide)
    add_header(slide, "Critique Workspace", "DRAFT")
    add_multiline(
        slide,
        "Prompts",
        [
            "What is the main contribution?",
            "Which evidence should a seminar audience remember?",
            "What is the sharpest discussion question?",
        ],
        0.78,
        1.65,
        5.6,
        2.4,
        fill=WHITE,
    )
    add_multiline(
        slide,
        "Placeholder verdict",
        ["Replace with concise review language.", "Keep slide readable for live discussion."],
        7.05,
        1.65,
        5.0,
        2.0,
        fill=MINT,
    )
    add_footer(slide)


def add_protected_slide(slide) -> None:
    add_background(slide, fill=WHITE)
    add_header(slide, "Lab Logistics", "PROTECTED")
    add_text(slide, "Do not edit: Transformer seminar logistics invariant", 0.82, 1.45, 7.5, 0.42, size=18, color=RED, bold=True)
    add_multiline(
        slide,
        "Seminar metadata",
        ["Reading group: SNUPI Office Agents", "Room: 302", "Time box: 25 min review + 20 min discussion"],
        0.82,
        2.18,
        5.8,
        1.9,
        fill=WHITE,
    )
    add_multiline(
        slide,
        "Benchmark invariant",
        ["This slide is used to detect collateral damage.", "Agent should preserve every text box on this slide."],
        7.05,
        2.18,
        5.0,
        1.9,
        fill=MIST,
    )
    add_footer(slide, "Protected context slide | Do not edit")


def build_seed(path: Path) -> None:
    prs = prs_16x9()
    for builder in [add_seed_slide1, add_seed_slide2, add_seed_slide3, add_seed_slide4, add_protected_slide]:
        builder(blank_slide(prs))
    prs.save(path)


def add_cover(slide, *, final_style=False) -> None:
    add_background(slide)
    add_box(slide, 0, 0, 13.333, 7.5, fill=PAPER, line=PAPER, radius=False, width=0)
    add_box(slide, 0.62, 0.55, 1.2, 5.95, fill=BLUE, line=BLUE, radius=False, width=0)
    add_text(slide, "Attention Is All You Need", 2.18, 1.08, 8.4, 0.64, size=34 if final_style else 32, color=BLUE if final_style else NAVY, bold=True, name="deck_title")
    add_text(slide, "Reading group review: Transformer architecture", 2.22, 1.86, 7.8, 0.36, size=18, color=INK, name="deck_subtitle")
    add_text(slide, "Vaswani et al. | arXiv:1706.03762 | NeurIPS 2017", 2.22, 2.36, 7.4, 0.26, size=12.5, color=STEEL)
    add_multiline(
        slide,
        "Main claim",
        [
            "The Transformer replaces recurrence with attention for sequence transduction.",
            "Parallelizable architecture plus attention layers delivered strong MT results.",
        ],
        2.22,
        3.05,
        5.35,
        1.55,
        fill=WHITE,
        name="main_claim_box",
    )
    add_multiline(
        slide,
        "Review lens",
        ["Architecture", "Attention formula", "Discussion risk"],
        8.05,
        3.05,
        3.65,
        1.55,
        fill=MINT,
        name="review_lens_box",
    )
    add_footer(slide)


def add_architecture_slide(slide, *, turn2=False, final=False) -> None:
    add_background(slide)
    add_header(slide, "Transformer Architecture", "FIGURE 1", final_style=final, name="architecture_title")
    if turn2 or final:
        add_box(slide, 0.62, 1.02, 5.85, 5.83, fill=WHITE, line=BLUE, width=1.4, name="architecture_stage")
        add_picture_fit(slide, FIGURE1_PATH, 1.52, 1.18, h=5.45, name="transformer_architecture_figure")
        add_text(slide, "Source figure: Transformer architecture (Figure 1)", 0.82, 6.55, 5.45, 0.26, size=10.5, color=STEEL, name="architecture_caption")
        add_multiline(
            slide,
            "Encoder stack",
            ["Self-attention + feed-forward layers", "Residual path and normalization keep depth stable."],
            7.02,
            1.30,
            5.0,
            1.25,
            fill=MIST,
            name="encoder_stack_callout",
        )
        add_multiline(
            slide,
            "Decoder stack",
            ["Masked attention prevents future-token leakage.", "Cross-attention reads encoder states."],
            7.02,
            2.85,
            5.0,
            1.25,
            fill=WHITE,
            name="decoder_stack_callout",
        )
        add_multiline(
            slide,
            "Multi-head attention bridge",
            ["The same attention primitive connects representation, decoding, and alignment."],
            7.02,
            4.40,
            5.0,
            1.25,
            fill=MINT,
            name="multi_head_bridge_callout",
        )
        add_pill(slide, "Turn 2 layout follows annotated screenshot", 7.05, 6.12, 3.6, fill=WHITE, line=GRAY)
    else:
        add_box(slide, 0.70, 1.18, 4.65, 5.55, fill=WHITE, line=BLUE, width=1.3, name="architecture_stage")
        add_picture_fit(slide, FIGURE1_PATH, 1.30, 1.34, h=4.95, name="transformer_architecture_figure")
        add_text(slide, "Source figure: Transformer architecture (Figure 1)", 0.88, 6.45, 4.15, 0.24, size=10.2, color=STEEL, name="architecture_caption")
        add_multiline(
            slide,
            "What to notice",
            [
                "Encoder and decoder are stacks, not recurrent chains.",
                "Attention is the central routing operation.",
                "The layout makes parallel training natural.",
            ],
            6.15,
            1.45,
            5.85,
            2.05,
            fill=MIST,
            name="architecture_summary",
        )
        add_multiline(
            slide,
            "Seminar note",
            ["Use this figure to explain why the paper changed the default mental model for sequence models."],
            6.15,
            4.10,
            5.85,
            1.15,
            fill=WHITE,
            name="architecture_note",
        )
    add_footer(slide)


def add_formula_slide(slide) -> None:
    add_background(slide)
    add_header(slide, "Scaled Dot-Product Attention", "FORMULA + FIGURE 2", name="formula_slide_title")
    add_box(slide, 0.62, 1.28, 6.35, 3.88, fill=WHITE, line=GRAY, width=1.2)
    add_picture_fit(slide, FIGURE2_PATH, 0.92, 1.55, w=5.75, name="attention_mechanism_figure")
    add_multiline(
        slide,
        "Scaled dot-product attention",
        [
            "Inputs are queries Q, keys K, and values V.",
            "Scaling by sqrt(d_k) stabilizes the softmax.",
        ],
        7.25,
        1.42,
        4.75,
        1.20,
        fill=MIST,
        name="formula_context",
    )
    add_text(
        slide,
        "Attention(Q, K, V) = softmax(QK^T / sqrt(d_k))V",
        7.18,
        2.98,
        5.05,
        0.42,
        size=17,
        color=NAVY,
        bold=True,
        name="attention_formula_text",
    )
    add_picture_fit(slide, FORMULA_IMAGE_PATH, 7.18, 3.60, w=4.65, name="attention_formula_image")
    add_text(slide, "Formula from paper PDF, Section 3.2.1", 7.22, 4.34, 4.45, 0.24, size=10.5, color=STEEL)
    add_text(slide, "Why it matters: attention replaces recurrence for sequence transduction", 0.72, 5.78, 8.8, 0.32, size=16, bold=True, color=INK, name="why_it_matters")
    add_footer(slide)


def add_verdict_slide(slide, *, final=False) -> None:
    add_background(slide)
    add_header(slide, "Paper Verdict", "DISCUSSION", final_style=final, name="verdict_title")
    add_text(slide, "Paper verdict", 0.78, 1.32, 4.2, 0.36, size=18, color=BLUE, bold=True)
    add_multiline(
        slide,
        "Strength: parallelizable sequence modeling",
        [
            "The paper made attention a practical replacement for recurrence.",
            "It offered a cleaner path to scale training across positions.",
        ],
        0.78,
        1.86,
        5.55,
        1.55,
        fill=WHITE,
        name="strength_card",
    )
    add_multiline(
        slide,
        "Question: what breaks when context grows?",
        [
            "Quadratic attention cost becomes the pressure point.",
            "Modern variants still negotiate this original design tradeoff.",
        ],
        6.95,
        1.86,
        5.2,
        1.55,
        fill=MIST,
        name="question_card",
    )
    add_multiline(
        slide,
        "For the reading group",
        ["Explain Figure 1 first.", "Use the formula to connect the architecture to computation.", "End with the scaling question."],
        0.78,
        4.22,
        5.55,
        1.65,
        fill=MINT,
        name="reading_group_card",
    )
    if final:
        add_text(
            slide,
            "Lab verdict: foundational, but evaluation now needs multimodal office tasks.",
            6.95,
            4.70,
            5.05,
            0.44,
            size=16,
            color=WHITE,
            bold=True,
            fill=BLUE,
            line=BLUE,
            margin=0.08,
            name="lab_verdict_footer",
        )
    add_footer(slide)


def build_deck(path: Path, *, turn2=False, final=False) -> None:
    prs = prs_16x9()
    add_cover(blank_slide(prs), final_style=final)
    add_architecture_slide(blank_slide(prs), turn2=turn2, final=final)
    add_formula_slide(blank_slide(prs))
    add_verdict_slide(blank_slide(prs), final=final)
    add_protected_slide(blank_slide(prs))
    prs.save(path)


def write_text_inputs() -> None:
    (INPUTS / "paper_source_notes.txt").write_text(
        "\n".join(
            [
                "Paper source: Attention Is All You Need",
                "arXiv:1706.03762, submitted 12 Jun 2017, version used by generator: PDF downloaded from https://arxiv.org/pdf/1706.03762",
                "Required formula string for slides: Attention(Q, K, V) = softmax(QK^T / sqrt(d_k))V",
                "Required Figure 1 caption: Source figure: Transformer architecture (Figure 1)",
                "Required Figure 2 caption/reference: Formula from paper PDF, Section 3.2.1",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    (INPUTS / "seminar_style_reference.svg").write_text(
        """<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <rect width="1280" height="720" fill="#faf8f2"/>
  <rect x="64" y="52" width="96" height="566" fill="#1f4e79"/>
  <text x="210" y="130" font-family="Arial" font-size="44" font-weight="700" fill="#1f4e79">SNUPI seminar style</text>
  <text x="210" y="188" font-family="Arial" font-size="24" fill="#56667a">Warm paper background, CUIF blue title, quiet source footers, and editable figure captions.</text>
  <rect x="210" y="270" width="420" height="110" rx="8" fill="#ffffff" stroke="#dae1e8" stroke-width="4"/>
  <text x="244" y="334" font-family="Arial" font-size="26" font-weight="700" fill="#172033">Evidence object</text>
  <rect x="690" y="270" width="360" height="110" rx="8" fill="#e5f7f0" stroke="#dae1e8" stroke-width="4"/>
  <text x="724" y="334" font-family="Arial" font-size="26" font-weight="700" fill="#172033">Discussion prompt</text>
  <text x="210" y="520" font-family="Arial" font-size="20" fill="#56667a">Final turn: style slide 1 title as #1F4E79, bold, 34 pt.</text>
</svg>
""",
        encoding="utf-8",
    )


def write_annotated_screenshot(seed_path: Path) -> None:
    target = INPUTS / "annotated_slide2_layout.png"
    base: Image.Image | None = None
    with tempfile.TemporaryDirectory(prefix="cuif_seed_render_") as tmp:
        preview = render_pptx_previews(seed_path, Path(tmp), "seed_for_annotation")
        images = preview.get("images", [])
        if len(images) >= 2:
            base = Image.open(images[1]).convert("RGB").resize((1280, 720))
    if base is None:
        base = Image.new("RGB", (1280, 720), "#faf8f2")
    overlay = Image.new("RGBA", base.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    try:
        font_big = ImageFont.truetype("Arial.ttf", 28)
        font_small = ImageFont.truetype("Arial.ttf", 22)
    except OSError:
        font_big = ImageFont.load_default()
        font_small = ImageFont.load_default()
    draw.rounded_rectangle((58, 96, 622, 652), radius=12, outline=(198, 61, 55, 255), width=7, fill=(255, 230, 80, 36))
    draw.text((88, 116), "ENLARGE FIGURE 1 HERE", fill=(198, 61, 55, 255), font=font_big)
    draw.line((1000, 172, 642, 230), fill=(198, 61, 55, 255), width=6)
    draw.polygon([(642, 230), (670, 212), (674, 244)], fill=(198, 61, 55, 255))
    draw.rounded_rectangle((676, 122, 1164, 260), radius=12, outline=(31, 78, 121, 255), width=5, fill=(229, 247, 240, 110))
    draw.text((704, 144), "Right rail: Encoder stack", fill=(31, 78, 121, 255), font=font_small)
    draw.rounded_rectangle((676, 286, 1164, 424), radius=12, outline=(31, 78, 121, 255), width=5, fill=(255, 255, 255, 120))
    draw.text((704, 308), "Right rail: Decoder stack", fill=(31, 78, 121, 255), font=font_small)
    draw.rounded_rectangle((676, 450, 1164, 588), radius=12, outline=(31, 78, 121, 255), width=5, fill=(229, 247, 240, 110))
    draw.text((704, 472), "Right rail: Multi-head attention bridge", fill=(31, 78, 121, 255), font=font_small)
    draw.text((84, 666), "Caption directly below figure", fill=(198, 61, 55, 255), font=font_small)
    Image.alpha_composite(base.convert("RGBA"), overlay).convert("RGB").save(target)


def main() -> None:
    ensure_dirs()
    download_pdf()
    crop_pdf_assets()
    write_text_inputs()
    seed = ARTIFACTS / "seed.pptx"
    build_seed(seed)
    write_annotated_screenshot(seed)
    build_deck(ARTIFACTS / "gold" / "turn1.pptx")
    build_deck(ARTIFACTS / "gold" / "turn2.pptx", turn2=True)
    build_deck(ARTIFACTS / "gold" / "final.pptx", turn2=True, final=True)
    shutil.copy2(ARTIFACTS / "gold" / "turn1.pptx", MOCK / "turn1" / "result.pptx")
    shutil.copy2(ARTIFACTS / "gold" / "turn2.pptx", MOCK / "turn2" / "result.pptx")
    shutil.copy2(ARTIFACTS / "gold" / "final.pptx", MOCK / "final" / "result.pptx")


if __name__ == "__main__":
    main()
