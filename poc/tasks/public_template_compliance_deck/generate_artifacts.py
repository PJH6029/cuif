from __future__ import annotations

import shutil
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

TASK_DIR = Path(__file__).resolve().parent
ARTIFACTS = TASK_DIR / "artifacts"
INPUTS = ARTIFACTS / "inputs"
GOLD = ARTIFACTS / "gold"
MOCK = TASK_DIR / "mock_outputs"

NAVY = RGBColor(17, 42, 74)       # #112A4A
BLUE = RGBColor(31, 78, 121)      # #1F4E79
TEAL = RGBColor(0, 111, 128)      # #006F80
GOLD_RGB = RGBColor(245, 197, 66) # #F5C542
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(248, 250, 252)
LIGHT_GRAY = RGBColor(223, 229, 236)
MID_GRAY = RGBColor(91, 107, 128)
PALE_BLUE = RGBColor(234, 243, 249)
PALE_GOLD = RGBColor(255, 248, 225)
PALE_TEAL = RGBColor(232, 247, 249)
RED = RGBColor(177, 54, 61)
GREEN = RGBColor(65, 134, 90)

SLIDE_W = 13.333333
SLIDE_H = 7.5
SEAL = INPUTS / "metro_civic_seal.png"


def ensure_dirs() -> None:
    for path in [ARTIFACTS, INPUTS, GOLD, MOCK / "turn1", MOCK / "final"]:
        path.mkdir(parents=True, exist_ok=True)


def prs_16x9() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def add_bg(slide, color=OFF_WHITE) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_shape(slide, x, y, w, h, *, fill=WHITE, line=LIGHT_GRAY, radius=False, name=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.0)
    if name:
        shape.name = name
    return shape


def add_text(slide, text, x, y, w, h, *, size=14, color=NAVY, bold=False, fill=None, line=None, margin=0.06, name=None):
    if fill is None and line is None:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        shape = add_shape(slide, x, y, w, h, fill=fill or WHITE, line=line or LIGHT_GRAY, radius=False, name=name)
    if name:
        shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_header_footer(slide, *, title: str, final_title: bool = False) -> None:
    add_bg(slide)
    add_shape(slide, 0, 0, SLIDE_W, 0.64, fill=NAVY, line=NAVY, name="template_header_band")
    slide.shapes.add_picture(str(SEAL), Inches(0.36), Inches(0.10), width=Inches(0.44), height=Inches(0.44))
    title_color = GOLD_RGB if final_title else WHITE
    title_size = 24 if final_title else 22
    add_text(slide, title, 0.94, 0.14, 7.5, 0.34, size=title_size, color=title_color, bold=True, name="template_title")
    add_text(slide, "METRO CIVIC GRANTS OFFICE", 9.05, 0.18, 3.6, 0.24, size=10.5, color=WHITE, bold=True, name="template_header_label")
    add_shape(slide, 0, 7.08, SLIDE_W, 0.42, fill=NAVY, line=NAVY, name="template_footer_band")
    add_text(
        slide,
        "PUBLIC TEMPLATE • Metro Civic Grants Office • Confidential draft",
        0.44,
        7.16,
        7.4,
        0.18,
        size=9,
        color=WHITE,
        bold=True,
        name="template_footer_text",
    )
    add_text(slide, "Template margin: 0.42 in minimum", 9.24, 7.16, 3.4, 0.18, size=8.5, color=WHITE, bold=False)


def add_section_band(slide, text: str, *, final_style: bool = False) -> None:
    add_shape(slide, 0.48, 0.88, 12.35, 0.42, fill=TEAL, line=TEAL, name="section_band")
    color = GOLD_RGB if final_style else WHITE
    add_text(slide, text, 0.64, 0.98, 6.1, 0.18, size=13, color=color, bold=True, name=text.lower().replace(" ", "_").replace("/", "_"))


def add_status_card(slide, heading: str, body: str, x: float, *, fill=WHITE):
    add_shape(slide, x, 1.56, 3.45, 1.34, fill=fill, line=RGBColor(185, 197, 211), radius=True, name=f"card_{heading.lower().replace(' ', '_')}")
    add_text(slide, heading, x + 0.20, 1.76, 2.9, 0.24, size=15.5, color=BLUE, bold=True, name=f"card_heading_{heading.lower().replace(' ', '_')}")
    add_text(slide, body, x + 0.20, 2.12, 2.95, 0.54, size=11.5, color=NAVY, bold=False)


def add_turn1_slide1(slide, *, final: bool = False) -> None:
    add_header_footer(slide, title="Metro Civic Grants: FY26 Oversight", final_title=final)
    add_section_band(slide, "PROGRAM STATUS BRIEFING")
    add_status_card(slide, "Program health", "On track: 14 active grants\n2 require technical assistance", 0.68, fill=WHITE)
    add_status_card(slide, "Funding status", "$42.8M obligated\n91% reporting complete", 4.94, fill=PALE_BLUE)
    add_status_card(slide, "Timeline", "Council packet due May 24\nAward letters by June 7", 9.20, fill=PALE_TEAL)
    add_shape(slide, 0.68, 3.42, 5.95, 1.92, fill=WHITE, line=RGBColor(179, 192, 207), radius=True, name="risk_watch_panel")
    add_text(slide, "Risk watch", 0.90, 3.66, 2.2, 0.26, size=15.5, color=BLUE, bold=True)
    add_text(
        slide,
        "Procurement timing is the main exposure.\nMitigation: publish pre-award checklist and hold weekly review.",
        0.92,
        4.08,
        5.25,
        0.76,
        size=11.5,
        color=NAVY,
    )
    add_shape(slide, 7.02, 3.42, 5.60, 1.92, fill=WHITE, line=RGBColor(179, 192, 207), radius=True, name="council_checkpoint_panel")
    add_text(slide, "Council checkpoint", 7.25, 3.66, 2.9, 0.26, size=15.5, color=BLUE, bold=True)
    add_text(slide, "Decision memo • Finance sign-off • District appendix", 7.26, 4.08, 4.80, 0.32, size=11.5, color=NAVY)
    add_text(slide, "Template rule: cards align to 0.42 in grid and stay clear of protected notices.", 7.26, 4.52, 4.95, 0.30, size=10.5, color=MID_GRAY)
    if final:
        add_shape(slide, 8.18, 5.56, 4.28, 0.92, fill=PALE_GOLD, line=GOLD_RGB, radius=True, name="oversight_callout")
        add_text(
            slide,
            "Oversight callout: add procurement review checkpoint by May 24.",
            8.36,
            5.82,
            3.78,
            0.28,
            size=12,
            color=NAVY,
            bold=True,
            name="oversight_callout_text",
        )
    add_text(
        slide,
        "Protected legal notice: FY26 allocations remain preliminary until council adoption.",
        0.70,
        6.42,
        7.75,
        0.24,
        size=9.5,
        color=MID_GRAY,
        bold=True,
        name="protected_legal_notice",
    )


def add_turn1_slide2(slide, *, final: bool = False) -> None:
    add_header_footer(slide, title="Metro Civic Grants: FY26 Oversight", final_title=False)
    add_section_band(slide, "COMPLIANCE MONITORING", final_style=final)
    add_shape(slide, 0.70, 1.48, 11.92, 3.05, fill=WHITE, line=RGBColor(179, 192, 207), radius=False, name="compliance_table")
    add_text(slide, "Compliance table", 0.92, 1.70, 2.30, 0.24, size=15, color=BLUE, bold=True, name="compliance_table_title")
    headers = [("Area", 0.96), ("Status", 3.88), ("Owner", 6.18), ("Next step", 8.48)]
    for header, x in headers:
        add_text(slide, header, x, 2.16, 1.75, 0.25, size=11.5, color=WHITE, bold=True, fill=TEAL, line=TEAL, margin=0.03)
    rows = [
        ("Reporting", "Green", "M. Chen", "Publish dashboard"),
        ("Procurement", "Amber", "R. Alvarez", "Review checklist"),
        ("Equity review", "Green", "T. Malik", "Archive scorecards"),
    ]
    y = 2.58
    for area, status, owner, step in rows:
        add_text(slide, area, 0.96, y, 2.42, 0.26, size=10.5, color=NAVY)
        add_text(slide, status, 3.88, y, 1.62, 0.26, size=10.5, color=GREEN if status == "Green" else RED, bold=True)
        add_text(slide, owner, 6.18, y, 1.62, 0.26, size=10.5, color=NAVY)
        add_text(slide, step, 8.48, y, 2.96, 0.26, size=10.5, color=NAVY)
        y += 0.56
    add_shape(slide, 0.70, 5.00, 5.48, 1.08, fill=PALE_BLUE, line=RGBColor(179, 192, 207), radius=True)
    add_text(slide, "Template note", 0.92, 5.20, 1.60, 0.22, size=12, color=BLUE, bold=True)
    add_text(slide, "Keep headers in teal, body text Aptos 10.5 pt, and table inside the grid.", 0.92, 5.54, 4.76, 0.26, size=10.5, color=NAVY)
    add_shape(slide, 7.02, 5.00, 5.60, 1.08, fill=WHITE, line=RGBColor(179, 192, 207), radius=True)
    add_text(slide, "Protected footer region", 7.25, 5.20, 2.45, 0.22, size=12, color=BLUE, bold=True)
    add_text(slide, "Footer and audit sentinel must not move into slide content.", 7.25, 5.54, 4.40, 0.26, size=10.5, color=NAVY)
    add_text(slide, "Do not edit: Metro Civic audit trail invariant.", 0.70, 6.42, 5.30, 0.24, size=9.5, color=MID_GRAY, bold=True, name="audit_trail_invariant")


def add_seed_slide1(slide) -> None:
    add_bg(slide, RGBColor(252, 252, 252))
    add_text(slide, "Metro Civic Grants rough notes", 0.58, 0.50, 7.2, 0.45, size=28, bold=True)
    add_text(slide, "Need public template applied using style_reference.svg and layout_grid.svg", 0.62, 1.12, 7.5, 0.32, size=13, color=MID_GRAY)
    add_text(slide, "Program health: On track; 14 active grants; 2 require technical assistance", 0.80, 1.82, 5.2, 0.55, size=16, fill=WHITE, line=LIGHT_GRAY)
    add_text(slide, "Funding status: $42.8M obligated; 91% reporting complete", 0.80, 2.60, 5.2, 0.55, size=16, fill=WHITE, line=LIGHT_GRAY)
    add_text(slide, "Timeline: Council packet due May 24; Award letters by June 7", 0.80, 3.38, 5.2, 0.55, size=16, fill=WHITE, line=LIGHT_GRAY)
    add_text(slide, "Risk watch: procurement timing is the main exposure", 6.88, 2.00, 5.0, 0.72, size=16, fill=WHITE, line=LIGHT_GRAY)
    add_text(slide, "Council checkpoint: Decision memo; Finance sign-off; District appendix", 6.88, 2.98, 5.0, 0.72, size=16, fill=WHITE, line=LIGHT_GRAY)
    add_text(slide, "Protected legal notice: FY26 allocations remain preliminary until council adoption.", 0.74, 6.24, 8.1, 0.32, size=10.5, color=MID_GRAY, bold=True)


def add_seed_slide2(slide) -> None:
    add_bg(slide, RGBColor(252, 252, 252))
    add_text(slide, "Compliance monitoring rough table", 0.58, 0.50, 7.2, 0.45, size=28, bold=True)
    add_text(slide, "Reporting Green M. Chen publish dashboard", 0.78, 1.62, 5.8, 0.45, size=15, fill=WHITE, line=LIGHT_GRAY)
    add_text(slide, "Procurement Amber R. Alvarez review checklist", 0.78, 2.28, 5.8, 0.45, size=15, fill=WHITE, line=LIGHT_GRAY)
    add_text(slide, "Equity review Green T. Malik archive scorecards", 0.78, 2.94, 5.8, 0.45, size=15, fill=WHITE, line=LIGHT_GRAY)
    add_text(slide, "Do not edit: Metro Civic audit trail invariant.", 0.74, 6.24, 5.8, 0.32, size=10.5, color=MID_GRAY, bold=True)


def create_seed() -> None:
    prs = prs_16x9()
    add_seed_slide1(prs.slides.add_slide(prs.slide_layouts[6]))
    add_seed_slide2(prs.slides.add_slide(prs.slide_layouts[6]))
    prs.save(ARTIFACTS / "seed.pptx")


def create_turn1() -> None:
    prs = prs_16x9()
    add_turn1_slide1(prs.slides.add_slide(prs.slide_layouts[6]), final=False)
    add_turn1_slide2(prs.slides.add_slide(prs.slide_layouts[6]), final=False)
    prs.save(GOLD / "turn1.pptx")


def create_final() -> None:
    prs = prs_16x9()
    add_turn1_slide1(prs.slides.add_slide(prs.slide_layouts[6]), final=True)
    add_turn1_slide2(prs.slides.add_slide(prs.slide_layouts[6]), final=True)
    prs.save(GOLD / "final.pptx")


def create_seal() -> None:
    image = Image.new("RGBA", (360, 360), (255, 255, 255, 0))
    draw = ImageDraw.Draw(image)
    draw.ellipse((18, 18, 342, 342), fill=(17, 42, 74, 255), outline=(245, 197, 66, 255), width=16)
    draw.ellipse((72, 72, 288, 288), fill=(255, 255, 255, 255), outline=(0, 111, 128, 255), width=8)
    try:
        font_big = ImageFont.truetype("DejaVuSans-Bold.ttf", 94)
        font_small = ImageFont.truetype("DejaVuSans-Bold.ttf", 25)
    except OSError:
        font_big = ImageFont.load_default()
        font_small = ImageFont.load_default()
    draw.text((180, 164), "MC", fill=(17, 42, 74, 255), font=font_big, anchor="mm")
    draw.text((180, 255), "GRANTS", fill=(0, 111, 128, 255), font=font_small, anchor="mm")
    image.save(SEAL)


def create_svg_assets() -> None:
    (INPUTS / "style_reference.svg").write_text(
        """<svg xmlns='http://www.w3.org/2000/svg' width='1280' height='720' viewBox='0 0 1280 720'>
  <rect width='1280' height='720' fill='#F8FAFC'/>
  <rect x='0' y='0' width='1280' height='62' fill='#112A4A'/>
  <rect x='0' y='680' width='1280' height='40' fill='#112A4A'/>
  <circle cx='54' cy='31' r='22' fill='#FFFFFF' stroke='#F5C542' stroke-width='5'/>
  <text x='92' y='39' font-family='Aptos, Arial' font-size='24' font-weight='700' fill='#FFFFFF'>Metro Civic Grants: FY26 Oversight</text>
  <rect x='46' y='84' width='1188' height='40' fill='#006F80'/>
  <text x='64' y='111' font-family='Aptos, Arial' font-size='18' font-weight='700' fill='#FFFFFF'>SECTION BAND: ALL CAPS WHITE / GOLD ON FINAL</text>
  <rect x='64' y='150' width='330' height='128' fill='#FFFFFF' stroke='#B9C5D3'/>
  <rect x='474' y='150' width='330' height='128' fill='#EAF3F9' stroke='#B9C5D3'/>
  <rect x='884' y='150' width='330' height='128' fill='#E8F7F9' stroke='#B9C5D3'/>
  <text x='64' y='338' font-family='Aptos, Arial' font-size='17' fill='#5B6B80'>Use navy #112A4A, teal #006F80, gold #F5C542, Aptos typography, and 0.42 inch margins.</text>
  <rect x='64' y='612' width='760' height='30' fill='none' stroke='#B1363D' stroke-dasharray='8 6'/>
  <text x='74' y='633' font-family='Aptos, Arial' font-size='15' font-weight='700' fill='#5B6B80'>Protected legal/audit notice region: do not edit or overlap.</text>
</svg>
""",
        encoding="utf-8",
    )
    (INPUTS / "layout_grid.svg").write_text(
        """<svg xmlns='http://www.w3.org/2000/svg' width='1280' height='720' viewBox='0 0 1280 720'>
  <rect width='1280' height='720' fill='#FFFFFF'/>
  <rect x='40' y='40' width='1200' height='640' fill='none' stroke='#112A4A' stroke-width='3'/>
  <line x1='40' y1='84' x2='1240' y2='84' stroke='#006F80' stroke-width='3'/>
  <line x1='40' y1='124' x2='1240' y2='124' stroke='#006F80' stroke-width='2'/>
  <rect x='65' y='150' width='330' height='128' fill='none' stroke='#1F4E79' stroke-width='3'/>
  <rect x='474' y='150' width='330' height='128' fill='none' stroke='#1F4E79' stroke-width='3'/>
  <rect x='884' y='150' width='330' height='128' fill='none' stroke='#1F4E79' stroke-width='3'/>
  <rect x='655' y='534' width='410' height='88' fill='none' stroke='#F5C542' stroke-width='4'/>
  <text x='70' y='192' font-family='Aptos, Arial' font-size='19' fill='#112A4A'>Three top cards align to this row</text>
  <text x='660' y='582' font-family='Aptos, Arial' font-size='19' fill='#112A4A'>Final callout must stay here</text>
  <text x='52' y='702' font-family='Aptos, Arial' font-size='14' fill='#5B6B80'>Normalized margins: x 0.04-0.96; footer reserved at bottom.</text>
</svg>
""",
        encoding="utf-8",
    )
    (INPUTS / "template_rules.txt").write_text(
        """Metro Civic public briefing template rules

1. Use the Metro Civic seal in the upper-left header of every slide.
2. Keep a navy #112A4A header and footer band on every slide.
3. Use a teal #006F80 all-caps section band below the header.
4. Use Aptos typography: title 22 pt white in turn 1, final title 24 pt gold #F5C542, section labels 13 pt bold.
5. Keep the three slide-1 status cards aligned across the top row inside the 0.42 inch margin grid.
6. Keep protected legal and audit notices in the bottom notice region; do not delete, restyle beyond template gray, or overlap them.
7. The final-turn oversight callout must stay in the approved lower-right grid zone and must not cover footer/protected text.
""",
        encoding="utf-8",
    )


def copy_mock_outputs() -> None:
    shutil.copy2(GOLD / "turn1.pptx", MOCK / "turn1" / "result.pptx")
    shutil.copy2(GOLD / "final.pptx", MOCK / "final" / "result.pptx")


def main() -> None:
    ensure_dirs()
    create_seal()
    create_svg_assets()
    create_seed()
    create_turn1()
    create_final()
    copy_mock_outputs()
    print(f"Generated public template compliance deck package at {TASK_DIR}")


if __name__ == "__main__":
    main()
