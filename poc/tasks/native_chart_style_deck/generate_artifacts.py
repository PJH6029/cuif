from __future__ import annotations

import csv
import html
import shutil
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

TASK_DIR = Path(__file__).resolve().parent
ARTIFACTS = TASK_DIR / "artifacts"
INPUTS = ARTIFACTS / "inputs"
GOLD = ARTIFACTS / "gold"
MOCK = TASK_DIR / "mock_outputs"

ROWS = [
    {"quarter": "Q1", "ticket_demand": 1180, "staffed_capacity": 1220, "gap": 40},
    {"quarter": "Q2", "ticket_demand": 1285, "staffed_capacity": 1320, "gap": 35},
    {"quarter": "Q3", "ticket_demand": 1420, "staffed_capacity": 1460, "gap": 40},
    {"quarter": "Q4", "ticket_demand": 1580, "staffed_capacity": 1510, "gap": -70},
]

NAVY = RGBColor(23, 32, 51)
BLUE = RGBColor(31, 78, 121)
DEEP_BLUE = RGBColor(31, 120, 180)
PALE_BLUE = RGBColor(166, 206, 227)
STEEL = RGBColor(93, 114, 135)
TEAL = RGBColor(42, 157, 143)
GOLD = RGBColor(232, 180, 66)
CREAM = RGBColor(250, 248, 242)
SKY = RGBColor(234, 244, 250)
MINT = RGBColor(232, 247, 240)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(229, 233, 238)
MID_GRAY = RGBColor(126, 137, 148)
RED = RGBColor(183, 76, 70)

# Rebind after color constant so paths retain clear names.
GOLD_DIR = ARTIFACTS / "gold"


def ensure_dirs() -> None:
    for path in [ARTIFACTS, INPUTS, GOLD_DIR, MOCK / "turn1", MOCK / "final"]:
        path.mkdir(parents=True, exist_ok=True)


def cell_xml(ref: str, value: Any, style: int | None = None) -> str:
    style_attr = f' s="{style}"' if style is not None else ""
    if isinstance(value, (int, float)):
        return f'<c r="{ref}"{style_attr}><v>{value}</v></c>'
    text = html.escape("" if value is None else str(value), quote=False)
    return f'<c r="{ref}" t="inlineStr"{style_attr}><is><t>{text}</t></is></c>'


def sheet_xml(rows: list[list[Any]], col_widths: list[float]) -> str:
    cols = "".join(f'<col min="{idx}" max="{idx}" width="{width}" customWidth="1"/>' for idx, width in enumerate(col_widths, start=1))
    row_xml = []
    for ridx, row in enumerate(rows, start=1):
        cells = []
        for cidx, value in enumerate(row, start=1):
            col = ""
            number = cidx
            while number:
                number, remainder = divmod(number - 1, 26)
                col = chr(65 + remainder) + col
            cells.append(cell_xml(f"{col}{ridx}", value, 1 if ridx == 1 else None))
        row_xml.append(f'<row r="{ridx}">{"".join(cells)}</row>')
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <cols>{cols}</cols>
  <sheetData>{"".join(row_xml)}</sheetData>
</worksheet>'''


def write_xlsx() -> None:
    data_rows = [
        ["Quarter", "Ticket demand", "Staffed capacity", "Capacity gap"],
        *[[row["quarter"], row["ticket_demand"], row["staffed_capacity"], row["gap"]] for row in ROWS],
    ]
    provenance = [
        ["Field", "Value"],
        ["Scenario", "Support operations capacity forecast"],
        ["Source owner", "Support Ops analytics extract"],
        ["Data cut", "FY2026 planning snapshot"],
        ["Required chart", "Native PowerPoint clustered column chart"],
        ["Negative gap meaning", "Demand exceeds staffed capacity"],
    ]
    with zipfile.ZipFile(INPUTS / "source_metrics.xlsx", "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>
  <Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>
  <Override PartName="/xl/worksheets/sheet2.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>
  <Override PartName="/xl/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.styles+xml"/>
  <Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
  <Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>
</Types>""")
        zf.writestr("_rels/.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>
  <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>
</Relationships>""")
        zf.writestr("xl/workbook.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <sheets>
    <sheet name="capacity_forecast" sheetId="1" r:id="rId1"/>
    <sheet name="provenance" sheetId="2" r:id="rId2"/>
  </sheets>
</workbook>""")
        zf.writestr("xl/_rels/workbook.xml.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet2.xml"/>
  <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/>
</Relationships>""")
        zf.writestr("xl/worksheets/sheet1.xml", sheet_xml(data_rows, [14, 18, 18, 15]))
        zf.writestr("xl/worksheets/sheet2.xml", sheet_xml(provenance, [24, 60]))
        zf.writestr("xl/styles.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<styleSheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <fonts count="2"><font><sz val="11"/><name val="Aptos"/></font><font><b/><sz val="11"/><name val="Aptos"/></font></fonts>
  <fills count="2"><fill><patternFill patternType="none"/></fill><fill><patternFill patternType="gray125"/></fill></fills>
  <borders count="1"><border><left/><right/><top/><bottom/><diagonal/></border></borders>
  <cellStyleXfs count="1"><xf numFmtId="0" fontId="0" fillId="0" borderId="0"/></cellStyleXfs>
  <cellXfs count="2"><xf numFmtId="0" fontId="0" fillId="0" borderId="0" xfId="0"/><xf numFmtId="0" fontId="1" fillId="0" borderId="0" xfId="0"/></cellXfs>
  <cellStyles count="1"><cellStyle name="Normal" xfId="0" builtinId="0"/></cellStyles>
</styleSheet>""")
        zf.writestr("docProps/app.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"><Application>CUIF PoC generator</Application></Properties>""")
        zf.writestr("docProps/core.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:title>Support capacity source metrics</dc:title><dc:creator>CUIF PoC generator</dc:creator></cp:coreProperties>""")


def write_source_files() -> None:
    with (INPUTS / "source_metrics.csv").open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=["quarter", "ticket_demand", "staffed_capacity", "gap"])
        writer.writeheader()
        writer.writerows(ROWS)
    (INPUTS / "analyst_notes.txt").write_text(
        "Support capacity deck notes\n\n"
        "Use the capacity_forecast sheet in source_metrics.xlsx as the source of truth.\n"
        "Create a native PowerPoint clustered column chart, not a pasted chart screenshot.\n"
        "Series must be named Ticket demand and Staffed capacity.\n"
        "Required audit strings include Q1 1180 1220 +40, Q2 1285 1320 +35, Q3 1420 1460 +40, and Q4 1580 1510 -70.\n"
        "Key narrative: Q4 demand exceeds staffed capacity by 70 tickets; add executive mitigation without changing the data.\n",
        encoding="utf-8",
    )


def write_svg_inputs() -> None:
    (INPUTS / "layout_reference.svg").write_text(
        """<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <rect width="1280" height="720" fill="#faf8f2"/>
  <rect x="58" y="36" width="720" height="88" rx="8" fill="#eaf4fa" stroke="#5d7287" stroke-width="3" stroke-dasharray="10 8"/>
  <text x="84" y="91" font-family="Arial" font-size="34" font-weight="700" fill="#172033">Title zone: Support Capacity Forecast</text>
  <rect x="62" y="142" width="780" height="425" rx="10" fill="#ffffff" stroke="#1f4e79" stroke-width="4"/>
  <text x="96" y="196" font-family="Arial" font-size="26" fill="#1f4e79">Native clustered column chart goes here</text>
  <rect x="64" y="590" width="780" height="72" rx="8" fill="#fff7df" stroke="#e8b442" stroke-width="3"/>
  <text x="92" y="638" font-family="Arial" font-size="23" fill="#172033">Insight strip under chart</text>
  <rect x="914" y="150" width="292" height="418" rx="10" fill="#e8f7f0" stroke="#2a9d8f" stroke-width="4"/>
  <text x="944" y="208" font-family="Arial" font-size="28" font-weight="700" fill="#172033">Right observation rail</text>
  <rect x="910" y="38" width="290" height="62" rx="26" fill="#ffffff" stroke="#e8b442" stroke-width="4"/>
  <text x="940" y="80" font-family="Arial" font-size="22" font-weight="700" fill="#172033">Badge area</text>
</svg>
""",
        encoding="utf-8",
    )
    (INPUTS / "style_reference.svg").write_text(
        """<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <rect width="1280" height="720" fill="#faf8f2"/>
  <text x="72" y="108" font-family="Arial" font-size="46" font-weight="700" fill="#1f4e79">CUIF operations briefing style</text>
  <text x="72" y="164" font-family="Arial" font-size="24" fill="#5d7287">Use cream canvas, CUIF-blue title, pale blue demand bars, deep blue capacity bars, and a teal right rail.</text>
  <rect x="72" y="236" width="220" height="72" rx="8" fill="#1f4e79"/>
  <text x="320" y="282" font-family="Arial" font-size="28" fill="#172033">Title blue #1F4E79, bold, 34 pt</text>
  <rect x="72" y="344" width="220" height="72" rx="8" fill="#a6cee3"/>
  <text x="320" y="390" font-family="Arial" font-size="28" fill="#172033">Ticket demand series #A6CEE3</text>
  <rect x="72" y="452" width="220" height="72" rx="8" fill="#1f78b4"/>
  <text x="320" y="498" font-family="Arial" font-size="28" fill="#172033">Staffed capacity series #1F78B4</text>
  <rect x="780" y="235" width="345" height="255" rx="10" fill="#e8f7f0" stroke="#2a9d8f" stroke-width="4"/>
  <text x="820" y="300" font-family="Arial" font-size="28" font-weight="700" fill="#172033">Right rail: concise mitigations</text>
  <line x1="820" y1="386" x2="1072" y2="386" stroke="#2a9d8f" stroke-width="4"/>
</svg>
""",
        encoding="utf-8",
    )


def make_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_fill(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color: RGBColor, width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_box(slide, x: float, y: float, w: float, h: float, *, fill: RGBColor = WHITE, line: RGBColor = LIGHT_GRAY, radius: bool = True, transparency: int = 0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill, transparency)
    set_line(shape, line, 1.2)
    return shape


def add_text(slide, text: str, x: float, y: float, w: float, h: float, *, size: float = 16, color: RGBColor = NAVY, bold: bool = False, fill: RGBColor | None = None, line: RGBColor | None = None, name: str | None = None, margin: float = 0.08):
    if fill is None:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        shape = add_box(slide, x, y, w, h, fill=fill, line=line or fill)
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    lines = text.split("\n")
    for idx, line_text in enumerate(lines):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = line_text
        font = run.font
        font.name = "Aptos"
        font.size = Pt(size)
        font.bold = bold
        font.color.rgb = color
    return shape


def add_background(slide) -> None:
    add_box(slide, 0, 0, 13.333, 7.5, fill=CREAM, line=CREAM, radius=False)
    add_box(slide, 0, 0, 13.333, 0.16, fill=BLUE, line=BLUE, radius=False)


def add_seed_title_slide(slide) -> None:
    add_background(slide)
    add_text(slide, "Support Capacity Briefing", 0.62, 0.48, 6.6, 0.54, size=30, bold=True)
    add_text(slide, "Seed state: rough planning notes need a native chart and executive layout.", 0.66, 1.14, 8.6, 0.34, size=14, color=STEEL)
    add_text(
        slide,
        "Raw notes\n- Use source_metrics.xlsx and source_metrics.csv\n- Show ticket demand vs staffed capacity\n- Highlight the Q4 capacity gap\n- Do not paste chart screenshots",
        0.72,
        1.95,
        5.5,
        2.45,
        size=18,
        fill=WHITE,
        line=LIGHT_GRAY,
    )
    add_text(slide, "Chart placeholder", 7.05, 2.0, 4.8, 2.25, size=22, bold=True, fill=SKY, line=LIGHT_GRAY)


def add_source_audit_slide(slide) -> None:
    add_background(slide)
    add_text(slide, "Workbook audit", 0.62, 0.46, 4.3, 0.52, size=30, bold=True)
    add_text(slide, "Source data: support capacity forecast", 0.66, 1.08, 5.4, 0.34, size=15, color=STEEL)
    audit = (
        "Q1 1180 1220 +40\n"
        "Q2 1285 1320 +35\n"
        "Q3 1420 1460 +40\n"
        "Q4 1580 1510 -70"
    )
    add_text(slide, audit, 0.78, 1.78, 4.7, 2.25, size=20, fill=WHITE, line=LIGHT_GRAY, name="capacity_audit_values")
    add_text(
        slide,
        "Fields: Quarter | Ticket demand | Staffed capacity | Capacity gap\nNegative gap means demand exceeds staffed capacity.\nSource workbook: source_metrics.xlsx / capacity_forecast",
        6.1,
        1.78,
        5.85,
        2.1,
        size=17,
        fill=MINT,
        line=TEAL,
    )


def add_protected_slide(slide) -> None:
    add_background(slide)
    add_text(slide, "Protected Benchmark Context", 0.62, 0.48, 5.6, 0.52, size=30, bold=True)
    add_text(
        slide,
        "Do not edit: native chart package invariant\nThis slide exists to detect collateral damage across turns.\nIt must remain text-identical to the seed deck.",
        0.78,
        1.72,
        8.3,
        1.55,
        size=18,
        fill=WHITE,
        line=LIGHT_GRAY,
    )
    add_text(slide, "CUIF focus: editable native chart + layout constraints", 0.8, 4.1, 6.9, 0.44, size=18, color=STEEL)


def chart_data() -> CategoryChartData:
    data = CategoryChartData()
    data.categories = [row["quarter"] for row in ROWS]
    data.add_series("Ticket demand", [row["ticket_demand"] for row in ROWS])
    data.add_series("Staffed capacity", [row["staffed_capacity"] for row in ROWS])
    return data


def add_capacity_chart(slide, x: float = 0.72, y: float = 1.52, w: float = 7.65, h: float = 4.35):
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), chart_data())
    graphic.name = "support_capacity_clustered_chart"
    chart = graphic.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Ticket demand vs staffed capacity"
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    for series_idx, series in enumerate(chart.series):
        series.has_data_labels = True
        series.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
        try:
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = PALE_BLUE if series_idx == 0 else DEEP_BLUE
        except Exception:
            pass
    try:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.minimum_scale = 1000
        chart.value_axis.maximum_scale = 1700
        chart.category_axis.tick_labels.font.size = Pt(10)
        chart.value_axis.tick_labels.font.size = Pt(9)
    except Exception:
        pass
    return graphic


def add_turn1_chart_slide(slide, *, final: bool = False) -> None:
    add_background(slide)
    title_size = 34 if final else 30
    title_color = BLUE if final else NAVY
    add_text(slide, "Support Capacity Forecast", 0.58, 0.36, 6.9, 0.56, size=title_size, color=title_color, bold=True, name="capacity_title")
    add_text(slide, "Ticket demand vs staffed capacity from source_metrics.xlsx", 0.62, 0.96, 7.3, 0.32, size=13.5, color=STEEL)
    add_capacity_chart(slide)
    add_text(slide, "Native clustered column chart from workbook data", 0.84, 5.92, 5.2, 0.34, size=13, color=STEEL)
    if final:
        add_text(slide, "Data cut: FY2026 forecast", 9.55, 0.38, 2.6, 0.54, size=14, bold=True, fill=WHITE, line=GOLD, name="data_cut_badge")
        add_text(
            slide,
            "Capacity observations\nLeader: Q4 demand 1580\nGap: Q4 -70 tickets\nMitigation: shift 2 weekend pods",
            9.38,
            1.58,
            3.12,
            2.25,
            size=16,
            fill=MINT,
            line=TEAL,
            name="capacity_observations_rail",
        )
        add_text(
            slide,
            "Insight: Q4 demand exceeds staffed capacity by 70 tickets; mitigation should shift two weekend pods before launch.",
            0.78,
            6.24,
            7.6,
            0.58,
            size=15.5,
            color=NAVY,
            bold=True,
            fill=RGBColor(255, 247, 223),
            line=GOLD,
            name="q4_capacity_gap_insight",
        )
    else:
        add_text(slide, "Chart caption: quarterly ticket demand and staffed support capacity", 0.78, 6.18, 7.3, 0.42, size=15, color=NAVY, fill=RGBColor(255, 247, 223), line=GOLD)


def build_seed() -> None:
    prs = make_prs()
    add_seed_title_slide(blank_slide(prs))
    add_source_audit_slide(blank_slide(prs))
    add_protected_slide(blank_slide(prs))
    prs.save(ARTIFACTS / "seed.pptx")


def build_turn1() -> None:
    prs = make_prs()
    add_turn1_chart_slide(blank_slide(prs), final=False)
    add_source_audit_slide(blank_slide(prs))
    add_protected_slide(blank_slide(prs))
    prs.save(GOLD_DIR / "turn1.pptx")


def build_final() -> None:
    prs = make_prs()
    add_turn1_chart_slide(blank_slide(prs), final=True)
    add_source_audit_slide(blank_slide(prs))
    add_protected_slide(blank_slide(prs))
    prs.save(GOLD_DIR / "final.pptx")


def main() -> None:
    ensure_dirs()
    write_source_files()
    write_xlsx()
    write_svg_inputs()
    build_seed()
    build_turn1()
    build_final()
    shutil.copy2(GOLD_DIR / "turn1.pptx", MOCK / "turn1" / "result.pptx")
    shutil.copy2(GOLD_DIR / "final.pptx", MOCK / "final" / "result.pptx")


if __name__ == "__main__":
    main()
