from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[4]
TASKS = ROOT / "poc" / "tasks"


def copy_seed_negative(task_id: str) -> None:
    task = TASKS / task_id
    source = task / "artifacts" / "seed.pptx"
    for turn in ("turn1", "final"):
        dest = task / "mock_outputs_negative_seed" / turn / "result.pptx"
        dest.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(source, dest)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, *, size: float = 12, color: RGBColor | None = None, bold: bool = False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or RGBColor(23, 32, 51)
    return shape


def screenshot_chart_deck(*, final: bool) -> Presentation:
    task = TASKS / "native_chart_style_deck"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image = task / "evidence" / "multimodal_review" / "rendered" / ("gold_final" if final else "gold_turn1") / "slide-1.png"
    if not image.exists():
        raise FileNotFoundError(f"required rendered screenshot missing: {image}")
    slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)
    # Overlay editable text in the intended regions so this fixture isolates the
    # native-chart/screenshot shortcut failure rather than merely failing OCR/text checks.
    if final:
        add_text(slide, "Support Capacity Forecast", 0.45, 0.32, 6.3, 0.55, size=34, color=RGBColor(31, 78, 121), bold=True)
        add_text(slide, "Data cut: FY2026 forecast", 9.45, 0.43, 2.8, 0.35, size=12, bold=True)
        add_text(slide, "Capacity observations\nLeader: Q4 demand 1580\nGap: Q4 -70 tickets\nMitigation: shift 2 weekend pods", 9.15, 1.55, 3.35, 1.35, size=13, bold=True)
        add_text(slide, "Insight: Q4 demand exceeds staffed capacity by 70 tickets; mitigation should shift two weekend pods before launch.", 0.75, 6.15, 7.7, 0.55, size=11)
    else:
        add_text(slide, "Support Capacity Forecast", 0.45, 0.32, 6.3, 0.55, size=32, bold=True)
        add_text(slide, "Ticket demand vs staffed capacity from source_metrics.xlsx", 0.50, 0.94, 6.7, 0.35, size=14)
        add_text(slide, "Chart caption: quarterly ticket demand and staffed support capacity", 0.75, 6.15, 7.0, 0.35, size=12)

    audit = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(
        audit,
        "Workbook audit\nSource data: support capacity forecast\nQ1 1180 1220 +40\nQ2 1285 1320 +35\nQ3 1420 1460 +40\nQ4 1580 1510 -70",
        0.65,
        0.65,
        8.5,
        3.5,
        size=20,
        bold=True,
    )
    protected = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(protected, "Do not edit: native chart package invariant", 0.75, 1.0, 8.0, 0.5, size=18, color=RGBColor(231, 111, 81), bold=True)
    return prs


def create_native_screenshot_negative() -> None:
    task = TASKS / "native_chart_style_deck"
    for turn, final in (("turn1", False), ("final", True)):
        dest = task / "mock_outputs_negative_screenshot" / turn / "result.pptx"
        dest.parent.mkdir(parents=True, exist_ok=True)
        screenshot_chart_deck(final=final).save(dest)


def main() -> None:
    copy_seed_negative("annotated_layout_repair_deck")
    copy_seed_negative("public_template_compliance_deck")
    create_native_screenshot_negative()
    print("created negative fixtures for A/B/C layout-constraint tasks")


if __name__ == "__main__":
    main()
