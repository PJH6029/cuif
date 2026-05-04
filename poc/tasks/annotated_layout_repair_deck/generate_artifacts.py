from __future__ import annotations

import shutil
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TASK_DIR = Path(__file__).resolve().parent
ARTIFACTS = TASK_DIR / "artifacts"
INPUTS = ARTIFACTS / "inputs"
GOLD = ARTIFACTS / "gold"
MOCK = TASK_DIR / "mock_outputs"

W = 10.0
H = 7.5

COLORS = {
    "navy": RGBColor(0x17, 0x20, 0x33),
    "blue": RGBColor(0x1F, 0x4E, 0x79),
    "teal": RGBColor(0x2A, 0x9D, 0x8F),
    "orange": RGBColor(0xE7, 0x6F, 0x51),
    "amber": RGBColor(0xF4, 0xA2, 0x61),
    "pale": RGBColor(0xF6, 0xF8, 0xFB),
    "line": RGBColor(0xD9, 0xE2, 0xEC),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "gray": RGBColor(0x61, 0x6E, 0x7C),
}


def ensure_dirs() -> None:
    for path in [INPUTS, GOLD, MOCK / "turn1", MOCK / "final"]:
        path.mkdir(parents=True, exist_ok=True)


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_textbox(slide, text: str, x: float, y: float, w: float, h: float, *, size: float = 16, color=COLORS["navy"], bold: bool = False, name: str | None = None, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return shape


def add_box(slide, text: str, x: float, y: float, w: float, h: float, *, fill=COLORS["white"], line=COLORS["line"], size: float = 15, color=COLORS["navy"], bold: bool = False, name: str | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.25)
    tf = shape.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.clear()
    for idx, line_text in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(size if idx == 0 else max(10, size - 3))
        run.font.color.rgb = color
        run.font.bold = bold if idx == 0 else False
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, *, color=COLORS["teal"], width: float = 2.0):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_common_protected_slides(prs: Presentation) -> None:
    slide1 = blank_slide(prs)
    add_textbox(slide1, "Atlas Reliability Review", 0.55, 0.55, 8.8, 0.6, size=30, bold=True, name="context_title")
    add_textbox(slide1, "Do not edit: Task A context invariant", 0.65, 1.35, 8.6, 0.35, size=14, color=COLORS["orange"], bold=True)
    add_box(slide1, "Context\nThis deck is used to validate visual-instruction grounding and non-target preservation for editable PPTX artifacts.", 0.65, 2.05, 8.4, 2.3, fill=COLORS["pale"], size=17, bold=True)
    add_textbox(slide1, "Protected slide: preserve exact text in all turns.", 0.65, 6.45, 8.6, 0.35, size=12, color=COLORS["gray"])


def add_slide3(prs: Presentation) -> None:
    slide3 = blank_slide(prs)
    add_textbox(slide3, "Runbook Evidence Appendix", 0.55, 0.55, 8.8, 0.5, size=28, bold=True)
    add_textbox(slide3, "Do not edit: runbook evidence invariant", 0.65, 1.25, 8.6, 0.35, size=14, color=COLORS["orange"], bold=True)
    add_box(slide3, "Protected evidence\n• Escalation policy version: 2026-Q2\n• On-call rotation owner: Platform SRE\n• Audit ticket: OPS-4721", 0.85, 2.0, 7.9, 2.5, fill=COLORS["pale"], size=17, bold=True)


def add_slide4(prs: Presentation) -> None:
    slide4 = blank_slide(prs)
    add_textbox(slide4, "Benchmark Controls", 0.55, 0.55, 8.8, 0.5, size=28, bold=True)
    add_textbox(slide4, "Do not edit: annotated repair benchmark invariant", 0.65, 1.25, 8.6, 0.35, size=14, color=COLORS["orange"], bold=True)
    add_box(slide4, "Scoring controls\nThis slide should remain unchanged while agents repair the annotated dashboard slide.", 0.85, 2.05, 7.9, 2.3, fill=COLORS["pale"], size=17, bold=True)


def add_seed_dashboard(prs: Presentation) -> None:
    slide = blank_slide(prs)
    add_textbox(slide, "Service Reliability Dashboard", 0.55, 0.35, 8.8, 0.55, size=28, bold=True, name="dashboard_title")
    add_box(slide, "Checkout SLA 97.2%\nTarget ≥ 99%", 6.4, 1.0, 2.6, 0.85, fill=RGBColor(0xFF, 0xFB, 0xF0), line=COLORS["amber"], size=15, bold=True, name="seed_checkout_metric")
    add_box(slide, "Queue backlog 42\nOwner gaps", 1.0, 2.0, 2.5, 0.85, fill=RGBColor(0xFF, 0xFB, 0xF0), line=COLORS["amber"], size=15, bold=True)
    add_box(slide, "Alert fatigue 18%\nDuplicate pages", 4.2, 4.8, 2.6, 0.85, fill=RGBColor(0xFF, 0xFB, 0xF0), line=COLORS["amber"], size=15, bold=True)
    add_box(slide, "Service health trend\nToo small and off-grid", 0.9, 3.25, 3.8, 1.3, fill=COLORS["pale"], line=COLORS["line"], size=14, bold=True, name="seed_trend_panel")
    add_line(slide, 1.2, 4.1, 1.9, 3.75, color=COLORS["teal"])
    add_line(slide, 1.9, 3.75, 2.6, 3.95, color=COLORS["teal"])
    add_line(slide, 2.6, 3.95, 3.4, 3.45, color=COLORS["teal"])
    add_line(slide, 3.4, 3.45, 4.2, 3.65, color=COLORS["teal"])
    add_box(slide, "Hotspot: checkout latency\nNeeds lower-right callout", 3.0, 2.05, 3.1, 1.05, fill=RGBColor(0xFF, 0xF2, 0xEC), line=COLORS["orange"], size=15, bold=True)
    add_textbox(slide, "Annotated screenshot says: top-row metrics, big trend panel, risk callout lower-right.", 0.7, 6.55, 8.8, 0.35, size=11, color=COLORS["gray"])


def add_repaired_dashboard(prs: Presentation, *, final: bool = False) -> None:
    slide = blank_slide(prs)
    add_textbox(slide, "Service Reliability Dashboard", 0.50, 0.28, 8.9, 0.55, size=32 if final else 29, color=COLORS["blue"] if final else COLORS["navy"], bold=True, name="dashboard_title")
    add_box(slide, "Checkout SLA 97.2%\nTarget ≥ 99%", 0.45, 1.05, 2.60, 0.90, fill=RGBColor(0xEB, 0xF4, 0xFF), line=COLORS["blue"], size=15, bold=True, name="checkout_metric_tile")
    add_box(slide, "Queue backlog 42\nOwner gaps", 3.70, 1.05, 2.60, 0.90, fill=RGBColor(0xEB, 0xF7, 0xF5), line=COLORS["teal"], size=15, bold=True, name="backlog_metric_tile")
    add_box(slide, "Alert fatigue 18%\nDuplicate pages", 6.95, 1.05, 2.60, 0.90, fill=RGBColor(0xFF, 0xF7, 0xED), line=COLORS["amber"], size=15, bold=True, name="fatigue_metric_tile")
    trend = add_box(slide, "Service health trend\nEditable trend panel: incidents fell after owner routing was restored.", 0.50, 2.25, 5.90, 3.25, fill=COLORS["pale"], line=COLORS["blue"], size=16, bold=True, name="service_health_trend_panel")
    trend.line.width = Pt(1.75)
    add_line(slide, 0.95, 4.95, 1.7, 4.55, color=COLORS["teal"], width=2.5)
    add_line(slide, 1.7, 4.55, 2.5, 4.75, color=COLORS["teal"], width=2.5)
    add_line(slide, 2.5, 4.75, 3.3, 4.30, color=COLORS["teal"], width=2.5)
    add_line(slide, 3.3, 4.30, 4.2, 4.55, color=COLORS["teal"], width=2.5)
    add_line(slide, 4.2, 4.55, 5.65, 4.05, color=COLORS["teal"], width=2.5)
    if final:
        add_box(slide, "Action owners\nSRE: Priya\nSupport: Marcus\nComms: Lina", 6.65, 2.40, 2.90, 1.55, fill=COLORS["white"], line=COLORS["teal"], size=15, bold=True, name="action_owners")
        add_textbox(slide, "Caption: checkout latency risk is concentrated in owner-routing gaps.", 0.65, 5.60, 5.60, 0.50, size=13, color=COLORS["gray"], name="trend_caption")
    else:
        add_box(slide, "Repair notes\nOnly slide 2 changed; protected slides remain unchanged.", 6.65, 2.40, 2.90, 1.35, fill=COLORS["white"], line=COLORS["line"], size=14, bold=True)
    add_box(slide, "Hotspot: checkout latency\nEscalation owner missing\nRoute rollback owner by Friday", 6.65, 4.75, 2.90, 1.65, fill=RGBColor(0xFF, 0xF2, 0xEC), line=COLORS["orange"], size=15, bold=True, name="checkout_latency_risk")


def make_seed() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    add_common_protected_slides(prs)
    add_seed_dashboard(prs)
    add_slide3(prs)
    add_slide4(prs)
    return prs


def make_turn1() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    add_common_protected_slides(prs)
    add_repaired_dashboard(prs, final=False)
    add_slide3(prs)
    add_slide4(prs)
    return prs


def make_final() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    add_common_protected_slides(prs)
    add_repaired_dashboard(prs, final=True)
    add_slide3(prs)
    add_slide4(prs)
    return prs


def write_repair_notes() -> None:
    (INPUTS / "repair_notes.txt").write_text(
        "Repair slide 2 only. Align Checkout SLA 97.2%, Queue backlog 42, and Alert fatigue 18% across the top row. "
        "Expand Service health trend into the large left-middle region. Move Hotspot: checkout latency into the lower-right. "
        "Preserve slides 1, 3, and 4 exactly, including all Do not edit sentinel strings.\n",
        encoding="utf-8",
    )


def write_reference_svg() -> None:
    svg = """<svg xmlns='http://www.w3.org/2000/svg' width='1000' height='750' viewBox='0 0 1000 750'>
  <rect width='1000' height='750' fill='#F7FAFC'/>
  <text x='50' y='65' font-family='Arial' font-size='34' font-weight='700' fill='#172033'>Annotated repair target for slide 2</text>
  <rect x='45' y='105' width='260' height='90' fill='#EBF4FF' stroke='#1F4E79' stroke-width='4'/>
  <rect x='370' y='105' width='260' height='90' fill='#EBF7F5' stroke='#2A9D8F' stroke-width='4'/>
  <rect x='695' y='105' width='260' height='90' fill='#FFF7ED' stroke='#F4A261' stroke-width='4'/>
  <text x='70' y='155' font-family='Arial' font-size='24' fill='#172033'>Checkout SLA metric</text>
  <text x='395' y='155' font-family='Arial' font-size='24' fill='#172033'>Queue backlog metric</text>
  <text x='720' y='155' font-family='Arial' font-size='24' fill='#172033'>Alert fatigue metric</text>
  <rect x='50' y='225' width='590' height='325' fill='#FFFFFF' stroke='#1F4E79' stroke-width='5'/>
  <text x='80' y='275' font-family='Arial' font-size='27' font-weight='700' fill='#172033'>Service health trend panel goes here</text>
  <polyline points='95,495 170,455 250,475 330,430 420,455 565,405' fill='none' stroke='#2A9D8F' stroke-width='8'/>
  <rect x='665' y='475' width='290' height='165' fill='#FFF2EC' stroke='#E76F51' stroke-width='5'/>
  <text x='690' y='535' font-family='Arial' font-size='25' font-weight='700' fill='#172033'>Lower-right risk callout</text>
  <path d='M445 675 C540 700 690 690 790 640' fill='none' stroke='#E76F51' stroke-width='5' stroke-dasharray='12 9'/>
  <text x='55' y='705' font-family='Arial' font-size='21' fill='#616E7C'>Do not alter slides 1, 3, or 4; preserve repaired layout during final turn.</text>
</svg>
"""
    (INPUTS / "annotated_repair_reference.svg").write_text(svg, encoding="utf-8")


def write_seed_screenshot() -> None:
    img = Image.new("RGB", (1000, 750), "#FFFFFF")
    draw = ImageDraw.Draw(img)
    try:
        font_big = ImageFont.truetype("DejaVuSans-Bold.ttf", 32)
        font_med = ImageFont.truetype("DejaVuSans.ttf", 22)
        font_small = ImageFont.truetype("DejaVuSans.ttf", 18)
    except Exception:
        font_big = font_med = font_small = None
    draw.text((55, 38), "Service Reliability Dashboard (messy source)", fill="#172033", font=font_big)
    draw.rounded_rectangle((640, 100, 900, 185), radius=12, fill="#FFFBF0", outline="#F4A261", width=3)
    draw.text((660, 126), "Checkout SLA 97.2%", fill="#172033", font=font_med)
    draw.rounded_rectangle((100, 200, 350, 285), radius=12, fill="#FFFBF0", outline="#F4A261", width=3)
    draw.text((120, 226), "Queue backlog 42", fill="#172033", font=font_med)
    draw.rounded_rectangle((420, 480, 680, 565), radius=12, fill="#FFFBF0", outline="#F4A261", width=3)
    draw.text((440, 506), "Alert fatigue 18%", fill="#172033", font=font_med)
    draw.rectangle((90, 325, 470, 455), fill="#F6F8FB", outline="#D9E2EC", width=3)
    draw.text((112, 350), "Service health trend", fill="#172033", font=font_med)
    draw.line((120, 410, 190, 375, 260, 395, 340, 345, 420, 365), fill="#2A9D8F", width=4)
    draw.rounded_rectangle((300, 205, 610, 310), radius=12, fill="#FFF2EC", outline="#E76F51", width=3)
    draw.text((320, 232), "Hotspot: checkout latency", fill="#172033", font=font_med)
    draw.text((70, 665), "Red annotations in the reference asset specify the repair target.", fill="#616E7C", font=font_small)
    img.save(INPUTS / "slide2_seed_screenshot.png")


def main() -> None:
    ensure_dirs()
    seed = make_seed()
    turn1 = make_turn1()
    final = make_final()
    seed.save(ARTIFACTS / "seed.pptx")
    turn1.save(GOLD / "turn1.pptx")
    final.save(GOLD / "final.pptx")
    shutil.copy2(GOLD / "turn1.pptx", MOCK / "turn1" / "result.pptx")
    shutil.copy2(GOLD / "final.pptx", MOCK / "final" / "result.pptx")
    write_repair_notes()
    write_reference_svg()
    write_seed_screenshot()


if __name__ == "__main__":
    main()
