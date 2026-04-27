from __future__ import annotations

import csv
import html
import json
import shutil
import urllib.request
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

TASK_DIR = Path(__file__).resolve().parent
ARTIFACTS = TASK_DIR / "artifacts"
INPUTS = ARTIFACTS / "inputs"
GOLD = ARTIFACTS / "gold"
MOCK = TASK_DIR / "mock_outputs"

DATA_URL = "https://api.worldbank.org/v2/country/USA;DEU;CHN;IND/indicator/EG.ELC.RNEW.ZS?format=json&per_page=20000"
METADATA_URL = "https://databank.worldbank.org/metadataglossary/world-development-indicators/series/EG.ELC.RNEW.ZS"
YEARS = ["2015", "2021"]
COUNTRY_ORDER = ["Germany", "India", "United States", "China"]
ISO_BY_COUNTRY = {"Germany": "DEU", "India": "IND", "United States": "USA", "China": "CHN"}
FALLBACK_VALUES = {
    ("Germany", "2015"): 30.0,
    ("Germany", "2021"): 39.8,
    ("India", "2015"): 14.9,
    ("India", "2021"): 19.1,
    ("United States", "2015"): 13.6,
    ("United States", "2021"): 20.3,
    ("China", "2015"): 23.9,
    ("China", "2021"): 28.4,
}
FALLBACK_CHANGES = {"Germany": 9.8, "India": 4.2, "United States": 6.6, "China": 4.6}

NAVY = RGBColor(23, 32, 51)
BLUE = RGBColor(31, 78, 121)
PALE_BLUE = RGBColor(166, 206, 227)
STEEL = RGBColor(93, 114, 135)
TEAL = RGBColor(42, 157, 143)
GOLDEN = RGBColor(232, 180, 66)
MINT = RGBColor(232, 247, 240)
SKY = RGBColor(234, 244, 250)
CREAM = RGBColor(250, 248, 242)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(229, 233, 238)
MID_GRAY = RGBColor(126, 137, 148)
RED = RGBColor(183, 76, 70)


def ensure_dirs() -> None:
    for path in [ARTIFACTS, INPUTS, GOLD, MOCK / "turn1", MOCK / "turn2", MOCK / "final"]:
        path.mkdir(parents=True, exist_ok=True)


def fetch_world_bank_payload() -> dict[str, Any]:
    request = urllib.request.Request(DATA_URL, headers={"User-Agent": "cuif-poc-task-generator/0.1"})
    try:
        with urllib.request.urlopen(request, timeout=20) as response:
            payload = json.loads(response.read().decode("utf-8"))
    except Exception as exc:
        return {"download_error": str(exc), "fallback_used": True, "source_url": DATA_URL, "rows": []}
    return {"download_error": None, "fallback_used": False, "source_url": DATA_URL, "payload": payload}


def selected_rows(source: dict[str, Any]) -> list[dict[str, Any]]:
    values = dict(FALLBACK_VALUES)
    precise_values = dict(FALLBACK_VALUES)
    payload = source.get("payload")
    if isinstance(payload, list) and len(payload) > 1 and isinstance(payload[1], list):
        for row in payload[1]:
            country = row.get("country", {}).get("value")
            year = str(row.get("date"))
            value = row.get("value")
            if country in COUNTRY_ORDER and year in YEARS and isinstance(value, (int, float)):
                precise_values[(country, year)] = float(value)
                values[(country, year)] = round(float(value), 1)

    rows = []
    for country in COUNTRY_ORDER:
        first = values[(country, "2015")]
        last = values[(country, "2021")]
        precise_first = precise_values[(country, "2015")]
        precise_last = precise_values[(country, "2021")]
        rows.append(
            {
                "country": country,
                "iso3": ISO_BY_COUNTRY[country],
                "value_2015": first,
                "value_2021": last,
                "change_pp": FALLBACK_CHANGES[country] if source.get("fallback_used") else round(precise_last - precise_first, 1),
            }
        )
    return rows


def write_csv(rows: list[dict[str, Any]]) -> None:
    with (INPUTS / "source_data.csv").open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=["country", "iso3", "value_2015", "value_2021", "change_pp"])
        writer.writeheader()
        writer.writerows(rows)


def cell_xml(ref: str, value: Any, style: int | None = None) -> str:
    style_attr = f' s="{style}"' if style is not None else ""
    if isinstance(value, (int, float)):
        return f'<c r="{ref}"{style_attr}><v>{value}</v></c>'
    text = html.escape("" if value is None else str(value), quote=False)
    return f'<c r="{ref}" t="inlineStr"{style_attr}><is><t>{text}</t></is></c>'


def sheet_xml(rows: list[list[Any]], col_widths: list[float]) -> str:
    cols = "".join(f'<col min="{idx}" max="{idx}" width="{width}" customWidth="1"/>' for idx, width in enumerate(col_widths, start=1))
    row_xml = []
    for ridx, row in enumerate(rows, start=1):
        cells = []
        for cidx, value in enumerate(row, start=1):
            col = ""
            number = cidx
            while number:
                number, remainder = divmod(number - 1, 26)
                col = chr(65 + remainder) + col
            cells.append(cell_xml(f"{col}{ridx}", value, 1 if ridx == 1 else None))
        row_xml.append(f'<row r="{ridx}">{"".join(cells)}</row>')
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <cols>{cols}</cols>
  <sheetData>{"".join(row_xml)}</sheetData>
</worksheet>'''


def write_xlsx(rows: list[dict[str, Any]]) -> None:
    selected = [
        ["Country", "ISO3", "2015 share", "2021 share", "Change pp"],
        *[[row["country"], row["iso3"], row["value_2015"], row["value_2021"], row["change_pp"]] for row in rows],
    ]
    provenance = [
        ["Field", "Value"],
        ["Indicator", "Renewable electricity output (% of total electricity output)"],
        ["World Bank code", "EG.ELC.RNEW.ZS"],
        ["Selected years", "2015 and 2021"],
        ["Source URL", DATA_URL],
        ["Metadata URL", METADATA_URL],
        ["Accessed", "2026-04-27"],
        ["Note", "2021 is the latest non-null year used for this PoC task extract."],
    ]
    with zipfile.ZipFile(INPUTS / "source_data.xlsx", "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>
  <Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>
  <Override PartName="/xl/worksheets/sheet2.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>
  <Override PartName="/xl/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.styles+xml"/>
  <Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
  <Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>
</Types>""")
        zf.writestr("_rels/.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>
  <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>
</Relationships>""")
        zf.writestr("xl/workbook.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <sheets>
    <sheet name="selected_countries" sheetId="1" r:id="rId1"/>
    <sheet name="provenance" sheetId="2" r:id="rId2"/>
  </sheets>
</workbook>""")
        zf.writestr("xl/_rels/workbook.xml.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet2.xml"/>
  <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/>
</Relationships>""")
        zf.writestr("xl/worksheets/sheet1.xml", sheet_xml(selected, [20, 10, 14, 14, 12]))
        zf.writestr("xl/worksheets/sheet2.xml", sheet_xml(provenance, [26, 110]))
        zf.writestr("xl/styles.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<styleSheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <fonts count="2"><font><sz val="11"/><name val="Aptos"/></font><font><b/><sz val="11"/><name val="Aptos"/></font></fonts>
  <fills count="2"><fill><patternFill patternType="none"/></fill><fill><patternFill patternType="gray125"/></fill></fills>
  <borders count="1"><border><left/><right/><top/><bottom/><diagonal/></border></borders>
  <cellStyleXfs count="1"><xf numFmtId="0" fontId="0" fillId="0" borderId="0"/></cellStyleXfs>
  <cellXfs count="2"><xf numFmtId="0" fontId="0" fillId="0" borderId="0" xfId="0"/><xf numFmtId="0" fontId="1" fillId="0" borderId="0" xfId="0"/></cellXfs>
  <cellStyles count="1"><cellStyle name="Normal" xfId="0" builtinId="0"/></cellStyles>
</styleSheet>""")
        zf.writestr("docProps/app.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"><Application>CUIF PoC generator</Application></Properties>""")
        zf.writestr("docProps/core.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/"><dc:title>Renewable electricity source data</dc:title><dc:creator>CUIF PoC generator</dc:creator></cp:coreProperties>""")


def write_source_notes(rows: list[dict[str, Any]], source: dict[str, Any]) -> None:
    lines = [
        "Source data for renewable_power_briefing_deck",
        "",
        "Indicator: Renewable electricity output (% of total electricity output)",
        "World Bank code: EG.ELC.RNEW.ZS",
        f"API: {DATA_URL}",
        f"Metadata: {METADATA_URL}",
        "Accessed: 2026-04-27",
        "Selected years: 2015 and 2021",
        "Note: 2021 is the latest non-null year used for this PoC task extract.",
        f"Fallback used during generation: {source.get('fallback_used')}",
        "",
        "Required rounded values:",
    ]
    for row in rows:
        lines.append(f"- {row['country']}: 2015={row['value_2015']}, 2021={row['value_2021']}, change={row['change_pp']} pts")
    (INPUTS / "source_notes.txt").write_text("\n".join(lines) + "\n", encoding="utf-8")


def write_svg_inputs() -> None:
    (INPUTS / "layout_reference.svg").write_text(
        """<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <rect width="1280" height="720" fill="#faf8f2"/>
  <rect x="64" y="44" width="770" height="110" rx="8" fill="#eaf4fa" stroke="#5d7287" stroke-width="3" stroke-dasharray="10 8"/>
  <text x="88" y="102" font-family="Arial" font-size="34" font-weight="700" fill="#172033">Title band: Renewable Electricity Snapshot</text>
  <rect x="72" y="188" width="800" height="392" rx="8" fill="#ffffff" stroke="#172033" stroke-width="4"/>
  <line x1="124" y1="526" x2="820" y2="526" stroke="#1f4e79" stroke-width="3" stroke-dasharray="7 8"/>
  <text x="122" y="235" font-family="Arial" font-size="28" fill="#1f4e79">Clustered column chart from source_data.xlsx</text>
  <text x="122" y="566" font-family="Arial" font-size="24" fill="#172033">Turn 2 insight sentence goes directly below chart</text>
  <rect x="928" y="188" width="286" height="392" rx="8" fill="#e8f7f0" stroke="#2a9d8f" stroke-width="4"/>
  <text x="954" y="240" font-family="Arial" font-size="26" font-weight="700" fill="#172033">Final right rail</text>
  <text x="954" y="286" font-family="Arial" font-size="22" fill="#172033">Leader</text>
  <text x="954" y="330" font-family="Arial" font-size="22" fill="#172033">Fastest gain</text>
  <text x="954" y="374" font-family="Arial" font-size="22" fill="#172033">Watch item</text>
  <rect x="958" y="48" width="254" height="66" rx="26" fill="#ffffff" stroke="#e8b442" stroke-width="4"/>
  <text x="986" y="91" font-family="Arial" font-size="22" font-weight="700" fill="#172033">Badge area</text>
</svg>
""",
        encoding="utf-8",
    )
    (INPUTS / "style_reference.svg").write_text(
        """<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <rect width="1280" height="720" fill="#faf8f2"/>
  <text x="72" y="110" font-family="Arial" font-size="46" font-weight="700" fill="#1f4e79">CUIF energy briefing style</text>
  <text x="72" y="164" font-family="Arial" font-size="24" fill="#5d7287">Use calm analyst typography, CUIF blue title, pale 2015 bars, and deep blue 2021 bars.</text>
  <rect x="72" y="236" width="220" height="74" rx="8" fill="#1f4e79"/>
  <text x="320" y="283" font-family="Arial" font-size="28" fill="#172033">Title blue #1F4E79</text>
  <rect x="72" y="344" width="220" height="74" rx="8" fill="#a6cee3"/>
  <text x="320" y="391" font-family="Arial" font-size="28" fill="#172033">2015 series #A6CEE3</text>
  <rect x="72" y="452" width="220" height="74" rx="8" fill="#1f78b4"/>
  <text x="320" y="499" font-family="Arial" font-size="28" fill="#172033">2021 series #1F78B4</text>
  <rect x="782" y="238" width="330" height="252" rx="10" fill="#e8f7f0" stroke="#2a9d8f" stroke-width="4"/>
  <text x="820" y="300" font-family="Arial" font-size="28" font-weight="700" fill="#172033">Right rail</text>
  <text x="820" y="350" font-family="Arial" font-size="22" fill="#172033">Short, source-backed observations</text>
  <line x1="820" y1="386" x2="1072" y2="386" stroke="#2a9d8f" stroke-width="4"/>
</svg>
""",
        encoding="utf-8",
    )


def set_font(run, size: float, color: RGBColor = NAVY, bold: bool = False) -> None:
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_text(slide, x, y, w, h, text, size=18, color=NAVY, bold=False, name: str | None = None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, line in enumerate(text.split("\n")):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        run = paragraph.add_run()
        run.text = line
        set_font(run, size, color, bold)
    return shape


def add_box(slide, x, y, w, h, fill=WHITE, line=LIGHT_GRAY, radius=False, name: str | None = None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.25)
    return shape


def add_chart(slide, rows: list[dict[str, Any]], x: float, y: float, w: float, h: float):
    chart_data = CategoryChartData()
    chart_data.categories = [row["country"] for row in rows]
    chart_data.add_series("2015", [row["value_2015"] for row in rows])
    chart_data.add_series("2021", [row["value_2021"] for row in rows])
    frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), chart_data)
    frame.name = "renewable_share_clustered_chart"
    chart = frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.maximum_scale = 45.0
    chart.value_axis.major_unit = 10.0
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.number_format = "0.0"
    chart.plots[0].data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    for series, color in zip(chart.series, [PALE_BLUE, BLUE], strict=True):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
    return frame


def add_source_footer(slide) -> None:
    add_text(slide, 0.72, 7.12, 11.9, 0.22, "Source: World Bank WDI EG.ELC.RNEW.ZS, accessed 2026-04-27. Values rounded to one decimal.", 7.5, MID_GRAY)


def add_protected_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0.65, 0.55, 12.0, 6.35, SKY, LIGHT_GRAY, radius=True)
    add_text(slide, 1.05, 1.05, 10.9, 0.55, "Protected benchmark context", 26, NAVY, True)
    add_text(slide, 1.05, 1.85, 10.8, 0.8, "Do not edit: renewable briefing invariant", 24, RED, True)
    add_text(
        slide,
        1.05,
        2.9,
        10.8,
        1.5,
        "This slide anchors collateral-damage checks for the CUIF PoC. It should survive every turn exactly.",
        18,
        STEEL,
    )
    return slide


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def build_seed(rows: list[dict[str, Any]]) -> None:
    prs = new_deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, 0.62, 0.42, 11.8, 0.56, "Renewables briefing workspace", 30, NAVY, True)
    add_text(slide, 0.64, 1.08, 10.4, 0.38, "Starter deck for an XLSX-to-PPT chart task", 15, STEEL)
    add_box(slide, 0.72, 1.72, 7.5, 4.55, WHITE, LIGHT_GRAY, radius=True)
    add_text(slide, 1.0, 2.0, 6.9, 0.38, "Chart placeholder", 22, BLUE, True)
    add_text(
        slide,
        1.0,
        2.55,
        6.85,
        1.45,
        "Build a native clustered column chart from artifacts/inputs/source_data.xlsx.\nUse countries Germany, India, United States, and China.\nCompare 2015 with 2021 renewable electricity share.",
        15,
        NAVY,
    )
    add_box(slide, 8.65, 1.72, 3.95, 4.55, MINT, TEAL, radius=True)
    add_text(slide, 8.95, 2.0, 3.35, 0.38, "Later-turn rail", 20, NAVY, True)
    add_text(slide, 8.95, 2.58, 3.25, 1.1, "Final turn will add observations here. Preserve data values and protected context.", 14, STEEL)
    add_source_footer(slide)

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide2, 0.64, 0.52, 11.8, 0.5, "Source workbook staging", 28, NAVY, True)
    add_text(slide2, 0.72, 1.28, 11.4, 0.56, "Use source_data.xlsx and source_notes.txt as immutable inputs. Do not modify the source workbook.", 16, STEEL)
    add_box(slide2, 0.72, 2.0, 11.7, 3.7, WHITE, LIGHT_GRAY, radius=True)
    add_text(slide2, 1.05, 2.35, 10.95, 0.42, "Turn 1 should replace this with a workbook audit slide.", 22, BLUE, True)

    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide3, 0.64, 0.52, 11.7, 0.52, "Interpretation notes draft", 28, NAVY, True)
    add_text(slide3, 0.72, 1.35, 11.4, 1.4, "Keep claims modest: the indicator is share of electricity output, not total energy use. Values are rounded to one decimal.", 18, STEEL)
    add_box(slide3, 0.72, 3.25, 11.7, 2.2, SKY, LIGHT_GRAY, radius=True)
    add_text(slide3, 1.05, 3.65, 10.95, 0.42, "Open slot for a concise insight after the chart is built.", 20, BLUE, True)

    add_protected_slide(prs)
    prs.save(ARTIFACTS / "seed.pptx")


def add_workbook_audit_slide(prs: Presentation, rows: list[dict[str, Any]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, 0.62, 0.42, 11.9, 0.52, "Workbook audit", 28, NAVY, True)
    add_text(slide, 0.66, 1.05, 11.5, 0.38, "Source data: renewable electricity output share", 15, STEEL)
    add_box(slide, 0.76, 1.65, 11.75, 4.56, WHITE, LIGHT_GRAY, radius=True)
    add_text(slide, 1.05, 1.92, 10.4, 0.3, "Country 2015 2021 Change", 13, BLUE, True)
    for idx, row in enumerate(rows):
        y = 2.48 + idx * 0.72
        row_text = f"{row['country']} {row['value_2015']:.1f} {row['value_2021']:.1f} +{row['change_pp']:.1f} pts"
        add_text(slide, 1.05, y, 10.4, 0.34, row_text, 14, NAVY, True)
    add_text(
        slide,
        1.05,
        5.62,
        10.8,
        0.36,
        "Rows validated from source_data.xlsx selected_countries sheet; chart must use the same rounded values.",
        12,
        STEEL,
    )
    add_source_footer(slide)


def add_interpretation_slide(prs: Presentation, final: bool = False) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, 0.62, 0.42, 11.9, 0.52, "Interpretation guardrails", 28, NAVY, True)
    add_box(slide, 0.78, 1.32, 11.7, 4.85, SKY, LIGHT_GRAY, radius=True)
    add_text(slide, 1.08, 1.7, 10.9, 0.4, "Do not overstate causality", 20, BLUE, True)
    add_text(
        slide,
        1.08,
        2.26,
        10.85,
        0.96,
        "The metric is renewable electricity share of total electricity output. It does not measure economy-wide energy use or emissions intensity.",
        17,
        NAVY,
    )
    add_text(slide, 1.08, 3.58, 10.9, 0.38, "Use a peer-comparison frame", 20, BLUE, True)
    tail = "Final deck should preserve the workbook audit values and cite World Bank WDI." if final else "Turn 2 should add one concise insight under the slide 1 chart."
    add_text(slide, 1.08, 4.14, 10.85, 0.85, tail, 17, NAVY)
    add_source_footer(slide)


def build_turn_deck(rows: list[dict[str, Any]], stage: str) -> Presentation:
    prs = new_deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_color = BLUE if stage == "final" else NAVY
    title_size = 34 if stage == "final" else 30
    add_text(slide, 0.64, 0.36, 8.9, 0.66, "Renewable Electricity Snapshot", title_size, title_color, True, name="slide1_title")
    add_text(slide, 0.68, 0.98, 8.0, 0.34, "2015 vs 2021 selected peers", 15, STEEL)
    if stage == "final":
        add_box(slide, 9.78, 0.48, 2.55, 0.44, WHITE, GOLDEN, radius=True, name="data_cut_badge")
        add_text(slide, 10.0, 0.55, 2.12, 0.24, "Data cut: 2015 vs 2021", 11, NAVY, True, name="data_cut_badge_text")
        chart_x, chart_y, chart_w, chart_h = 0.72, 1.58, 8.05, 4.62
    else:
        chart_x, chart_y, chart_w, chart_h = 0.78, 1.58, 10.78, 4.36

    add_text(slide, chart_x, 1.25, chart_w, 0.28, "Renewable share of electricity output (%)", 13, BLUE, True)
    add_chart(slide, rows, chart_x, chart_y, chart_w, chart_h)

    if stage in {"turn2", "final"}:
        insight_y = 6.13 if stage == "turn2" else 6.24
        add_text(
            slide,
            chart_x,
            insight_y,
            chart_w,
            0.42,
            "Insight: Germany leads at 39.8% in 2021; the U.S. closes part of the gap with a +6.6 pt gain.",
            14,
            NAVY,
            True,
            name="chart_insight",
        )
        add_text(slide, chart_x, insight_y + 0.42, chart_w, 0.28, "Source: World Bank WDI EG.ELC.RNEW.ZS; latest non-null year in workbook is 2021.", 9, MID_GRAY)
    else:
        add_text(slide, chart_x, 6.12, chart_w, 0.32, "Insight pending after chart review.", 12, MID_GRAY)

    if stage == "final":
        add_box(slide, 9.18, 1.58, 3.28, 4.62, MINT, TEAL, radius=True, name="observation_rail")
        add_text(slide, 9.48, 1.92, 2.7, 0.34, "Key observations", 17, NAVY, True)
        add_text(slide, 9.48, 2.58, 2.7, 0.36, "Leader: Germany 39.8%", 15, NAVY, True)
        add_text(slide, 9.48, 3.36, 2.7, 0.36, "Fastest gain: Germany +9.8 pts", 15, NAVY, True)
        add_text(slide, 9.48, 4.14, 2.7, 0.36, "Watch item: U.S. 20.3%", 15, NAVY, True)
        add_text(slide, 9.48, 5.0, 2.55, 0.52, "Keep source workbook unchanged.", 12, STEEL)

    add_source_footer(slide)
    add_workbook_audit_slide(prs, rows)
    add_interpretation_slide(prs, final=stage == "final")
    add_protected_slide(prs)
    return prs


def write_manifest() -> None:
    (TASK_DIR / "manifest.yaml").write_text(
        """manifest_version: "0.1"
id: renewable_power_briefing_deck
title: Renewable power briefing from source workbook
primary_artifact_family: pptx
artifact_families: [pptx, xlsx]
tracks: [open_tool, gui]
scoring:
  policy: sum_points
judge:
  api_key_env: OPENAI_API_KEY
artifacts:
  package:
    seed:
      path: artifacts/seed.pptx
      type: pptx
      role: seed
    source_data:
      path: artifacts/inputs/source_data.xlsx
      type: xlsx
      role: source_input
    source_notes:
      path: artifacts/inputs/source_notes.txt
      type: txt
      role: source_input
    source_csv:
      path: artifacts/inputs/source_data.csv
      type: txt
      role: source_input
    raw_world_bank_json:
      path: artifacts/inputs/world_bank_renewable_api.json
      type: json
      role: source_input
    layout_reference:
      path: artifacts/inputs/layout_reference.svg
      type: svg
      role: instruction_input
    style_reference:
      path: artifacts/inputs/style_reference.svg
      type: svg
      role: style_input
    gold_turn1:
      path: artifacts/gold/turn1.pptx
      type: pptx
      role: gold
    gold_turn2:
      path: artifacts/gold/turn2.pptx
      type: pptx
      role: gold
    gold_final:
      path: artifacts/gold/final.pptx
      type: pptx
      role: gold
  expected_outputs:
    turn1:
      result:
        path: result.pptx
        type: pptx
    turn2:
      result:
        path: result.pptx
        type: pptx
    final:
      result:
        path: result.pptx
        type: pptx
turns:
  - id: turn1
    instruction: >-
      Using the seed deck, source_data.xlsx, source_notes.txt, and
      layout_reference.svg, create a four-slide renewable-electricity briefing.
      Slide 1 should follow the layout reference: add the title "Renewable
      Electricity Snapshot" near the top, the subtitle "2015 vs 2021 selected
      peers", and a native clustered column chart named
      "renewable_share_clustered_chart" using the workbook values for Germany,
      India, United States, and China. The chart must compare the series "2015"
      and "2021" with the rounded values in source_data.xlsx. Add the chart
      caption "Renewable share of electricity output (%)". Slide 2 should be a
      workbook audit containing the exact strings "Workbook audit", "Source
      data: renewable electricity output share", "Germany 30.0 39.8 +9.8 pts",
      "India 14.9 19.1 +4.2 pts", "United States 13.6 20.3 +6.6 pts", and
      "China 23.9 28.4 +4.6 pts". Slide 3 should explain that the metric is
      electricity-output share, not economy-wide energy use. Preserve slide 4
      exactly because it contains the sentinel "Do not edit: renewable briefing
      invariant".
    expected_output: result
    checks:
      - id: renewable_t1_file_exists
        evaluator: file_exists
        artifact: run.outputs.turn1.result
        points: 1
      - id: renewable_t1_slide_count
        evaluator: pptx_slide_count
        artifact: run.outputs.turn1.result
        points: 1
        depends_on: [renewable_t1_file_exists]
        params: {count: 4}
      - id: renewable_t1_core_texts
        evaluator: pptx_text_contains
        artifact: run.outputs.turn1.result
        points: 3
        depends_on: [renewable_t1_file_exists]
        params:
          texts:
            - "Renewable Electricity Snapshot"
            - "2015 vs 2021 selected peers"
            - "Renewable share of electricity output (%)"
            - "World Bank WDI EG.ELC.RNEW.ZS"
      - id: renewable_t1_chart_data
        evaluator: pptx_chart_data
        artifact: run.outputs.turn1.result
        points: 5
        depends_on: [renewable_t1_file_exists]
        params:
          selector: {slide: 1, name: "renewable_share_clustered_chart"}
          chart_type: "COLUMN_CLUSTERED"
          categories: ["Germany", "India", "United States", "China"]
          series:
            - {name: "2015", values: [30.0, 14.9, 13.6, 23.9]}
            - {name: "2021", values: [39.8, 19.1, 20.3, 28.4]}
          value_tolerance: 0.15
      - id: renewable_t1_title_region
        evaluator: pptx_bbox_region
        artifact: run.outputs.turn1.result
        points: 2
        depends_on: [renewable_t1_core_texts]
        params:
          selector: {slide: 1, text_contains: "Renewable Electricity Snapshot"}
          region: {slide: 1, x_min: 0.04, y_min: 0.03, x_max: 0.72, y_max: 0.18}
          tolerance: 0.02
      - id: renewable_t1_chart_region
        evaluator: pptx_bbox_region
        artifact: run.outputs.turn1.result
        points: 2
        depends_on: [renewable_t1_chart_data]
        params:
          selector: {slide: 1, name: "renewable_share_clustered_chart"}
          region: {slide: 1, x_min: 0.05, y_min: 0.20, x_max: 0.90, y_max: 0.84}
          tolerance: 0.03
      - id: renewable_t1_audit_values
        evaluator: pptx_text_contains
        artifact: run.outputs.turn1.result
        points: 4
        depends_on: [renewable_t1_file_exists]
        params:
          texts:
            - "Workbook audit"
            - "Source data: renewable electricity output share"
            - "Germany 30.0 39.8 +9.8 pts"
            - "India 14.9 19.1 +4.2 pts"
            - "United States 13.6 20.3 +6.6 pts"
            - "China 23.9 28.4 +4.6 pts"
      - id: renewable_t1_guardrail_text
        evaluator: pptx_text_contains
        artifact: run.outputs.turn1.result
        points: 2
        depends_on: [renewable_t1_file_exists]
        params:
          text: "The metric is renewable electricity share of total electricity output."
      - id: renewable_t1_protected_slide_preserved
        evaluator: pptx_preservation_diff
        artifact: run.outputs.turn1.result
        points: 3
        depends_on: [renewable_t1_slide_count]
        params:
          reference: package.seed
          compare: exact_slide_text
          slides: [4]
          protected_texts: ["Do not edit: renewable briefing invariant"]
      - id: renewable_t1_rendered_review
        evaluator: rendered_layout_review
        artifact: run.outputs.turn1.result
        points: 0
        optional: true
        diagnostic: true
        params:
          note: Optional rendered review for chart placement against layout_reference.svg.
  - id: turn2
    instruction: >-
      Revise the turn 1 deck by adding a one-sentence insight directly below
      the chart on slide 1: "Insight: Germany leads at 39.8% in 2021; the U.S.
      closes part of the gap with a +6.6 pt gain." Add the footnote "Source:
      World Bank WDI EG.ELC.RNEW.ZS; latest non-null year in workbook is
      2021." Keep the native chart data and slide 2 workbook audit exactly as
      they were after turn 1, and preserve slide 4 exactly.
    expected_output: result
    checks:
      - id: renewable_t2_file_exists
        evaluator: file_exists
        artifact: run.outputs.turn2.result
        points: 1
      - id: renewable_t2_slide_count
        evaluator: pptx_slide_count
        artifact: run.outputs.turn2.result
        points: 1
        depends_on: [renewable_t2_file_exists]
        params: {count: 4}
      - id: renewable_t2_chart_data_preserved
        evaluator: pptx_chart_data
        artifact: run.outputs.turn2.result
        points: 4
        depends_on: [renewable_t2_file_exists]
        params:
          selector: {slide: 1, name: "renewable_share_clustered_chart"}
          chart_type: "COLUMN_CLUSTERED"
          categories: ["Germany", "India", "United States", "China"]
          series:
            - {name: "2015", values: [30.0, 14.9, 13.6, 23.9]}
            - {name: "2021", values: [39.8, 19.1, 20.3, 28.4]}
          value_tolerance: 0.15
      - id: renewable_t2_insight_present
        evaluator: pptx_text_contains
        artifact: run.outputs.turn2.result
        points: 3
        depends_on: [renewable_t2_file_exists]
        params:
          text: "Insight: Germany leads at 39.8% in 2021; the U.S. closes part of the gap with a +6.6 pt gain."
      - id: renewable_t2_insight_region
        evaluator: pptx_bbox_region
        artifact: run.outputs.turn2.result
        points: 2
        depends_on: [renewable_t2_insight_present]
        params:
          selector: {slide: 1, text_contains: "Insight: Germany leads at 39.8%"}
          region: {slide: 1, x_min: 0.05, y_min: 0.78, x_max: 0.90, y_max: 0.92}
          tolerance: 0.03
      - id: renewable_t2_footnote_present
        evaluator: pptx_text_contains
        artifact: run.outputs.turn2.result
        points: 2
        depends_on: [renewable_t2_file_exists]
        params:
          text: "Source: World Bank WDI EG.ELC.RNEW.ZS; latest non-null year in workbook is 2021."
      - id: renewable_t2_slide2_preserved
        evaluator: pptx_preservation_diff
        artifact: run.outputs.turn2.result
        points: 3
        depends_on: [renewable_t2_slide_count]
        params:
          reference: run.outputs.turn1.result
          compare: exact_slide_text
          slides: [2]
          protected_texts:
            - "Workbook audit"
            - "Germany 30.0 39.8 +9.8 pts"
            - "United States 13.6 20.3 +6.6 pts"
      - id: renewable_t2_slide4_preserved
        evaluator: pptx_preservation_diff
        artifact: run.outputs.turn2.result
        points: 3
        depends_on: [renewable_t2_slide_count]
        params:
          reference: run.outputs.turn1.result
          compare: exact_slide_text
          slides: [4]
          protected_texts: ["Do not edit: renewable briefing invariant"]
      - id: renewable_t2_rendered_review
        evaluator: rendered_layout_review
        artifact: run.outputs.turn2.result
        points: 0
        optional: true
        diagnostic: true
        params:
          note: Optional rendered review for insight placement and chart preservation.
  - id: final
    instruction: >-
      Finalize the deck using style_reference.svg. On slide 1, apply the CUIF
      title style to "Renewable Electricity Snapshot": blue #1F4E79, bold, 34
      pt. Keep the source chart data unchanged, move/size the chart into the
      left chart region, and add a top-right badge with "Data cut: 2015 vs
      2021". Add a right-side observations rail containing exactly "Leader:
      Germany 39.8%", "Fastest gain: Germany +9.8 pts", and "Watch item: U.S.
      20.3%". Keep the turn 2 insight sentence and source footnote. Preserve
      slide 2 exactly from turn 2 and preserve slide 4 exactly.
    expected_output: result
    checks:
      - id: renewable_final_file_exists
        evaluator: file_exists
        artifact: run.outputs.final.result
        points: 1
      - id: renewable_final_slide_count
        evaluator: pptx_slide_count
        artifact: run.outputs.final.result
        points: 1
        depends_on: [renewable_final_file_exists]
        params: {count: 4}
      - id: renewable_final_title_present
        evaluator: pptx_text_contains
        artifact: run.outputs.final.result
        points: 1
        depends_on: [renewable_final_file_exists]
        params: {text: "Renewable Electricity Snapshot"}
      - id: renewable_final_title_style
        evaluator: pptx_style_check
        artifact: run.outputs.final.result
        points: 4
        depends_on: [renewable_final_title_present]
        params:
          selector: {slide: 1, text_contains: "Renewable Electricity Snapshot"}
          font_color: "#1F4E79"
          font_size_pt: 34
          font_size_tolerance: 0.25
          bold: true
      - id: renewable_final_chart_data_preserved
        evaluator: pptx_chart_data
        artifact: run.outputs.final.result
        points: 4
        depends_on: [renewable_final_file_exists]
        params:
          selector: {slide: 1, name: "renewable_share_clustered_chart"}
          chart_type: "COLUMN_CLUSTERED"
          categories: ["Germany", "India", "United States", "China"]
          series:
            - {name: "2015", values: [30.0, 14.9, 13.6, 23.9]}
            - {name: "2021", values: [39.8, 19.1, 20.3, 28.4]}
          value_tolerance: 0.15
      - id: renewable_final_chart_left_region
        evaluator: pptx_bbox_region
        artifact: run.outputs.final.result
        points: 2
        depends_on: [renewable_final_chart_data_preserved]
        params:
          selector: {slide: 1, name: "renewable_share_clustered_chart"}
          region: {slide: 1, x_min: 0.04, y_min: 0.18, x_max: 0.70, y_max: 0.86}
          tolerance: 0.03
      - id: renewable_final_observation_texts
        evaluator: pptx_text_contains
        artifact: run.outputs.final.result
        points: 4
        depends_on: [renewable_final_file_exists]
        params:
          texts:
            - "Data cut: 2015 vs 2021"
            - "Leader: Germany 39.8%"
            - "Fastest gain: Germany +9.8 pts"
            - "Watch item: U.S. 20.3%"
      - id: renewable_final_badge_region
        evaluator: pptx_bbox_region
        artifact: run.outputs.final.result
        points: 2
        depends_on: [renewable_final_observation_texts]
        params:
          selector: {slide: 1, text_contains: "Data cut: 2015 vs 2021"}
          region: {slide: 1, x_min: 0.72, y_min: 0.04, x_max: 0.96, y_max: 0.16}
          tolerance: 0.03
      - id: renewable_final_insight_preserved
        evaluator: pptx_text_contains
        artifact: run.outputs.final.result
        points: 2
        depends_on: [renewable_final_file_exists]
        params:
          texts:
            - "Insight: Germany leads at 39.8% in 2021; the U.S. closes part of the gap with a +6.6 pt gain."
            - "Source: World Bank WDI EG.ELC.RNEW.ZS; latest non-null year in workbook is 2021."
      - id: renewable_final_slide2_preserved
        evaluator: pptx_preservation_diff
        artifact: run.outputs.final.result
        points: 3
        depends_on: [renewable_final_slide_count]
        params:
          reference: run.outputs.turn2.result
          compare: exact_slide_text
          slides: [2]
          protected_texts:
            - "Workbook audit"
            - "Germany 30.0 39.8 +9.8 pts"
            - "China 23.9 28.4 +4.6 pts"
      - id: renewable_final_slide4_preserved
        evaluator: pptx_preservation_diff
        artifact: run.outputs.final.result
        points: 3
        depends_on: [renewable_final_slide_count]
        params:
          reference: run.outputs.turn2.result
          compare: exact_slide_text
          slides: [4]
          protected_texts: ["Do not edit: renewable briefing invariant"]
      - id: optional_llm_renewable_insight_judge
        evaluator: llm_text_rubric
        artifact: run.outputs.final.result
        points: 0
        optional: true
        diagnostic: true
        params:
          prompt: "Check whether the final deck insight follows from the renewable electricity values in the slide text/chart."
          rubric: "Pass if the deck correctly says Germany leads at 39.8% in 2021, Germany has the largest gain at +9.8 points, and the U.S. value is 20.3% without overclaiming beyond electricity-output share."
      - id: optional_vlm_renewable_layout_judge
        evaluator: vlm_layout_rubric
        artifact: run.outputs.final.result
        points: 0
        optional: true
        diagnostic: true
        params:
          prompt: "Inspect the rendered final renewable electricity deck."
          rubric: "Pass if slide 1 has a readable chart on the left, a clear observations rail on the right, a top-right badge, and no obvious overlap or collateral damage."
""",
        encoding="utf-8",
    )


def main() -> None:
    ensure_dirs()
    source = fetch_world_bank_payload()
    (INPUTS / "world_bank_renewable_api.json").write_text(json.dumps(source, indent=2, sort_keys=True), encoding="utf-8")
    rows = selected_rows(source)
    write_csv(rows)
    write_xlsx(rows)
    write_source_notes(rows, source)
    write_svg_inputs()
    build_seed(rows)
    for stage, filename in [("turn1", "turn1.pptx"), ("turn2", "turn2.pptx"), ("final", "final.pptx")]:
        deck = build_turn_deck(rows, stage)
        deck.save(GOLD / filename)
    shutil.copy2(GOLD / "turn1.pptx", MOCK / "turn1" / "result.pptx")
    shutil.copy2(GOLD / "turn2.pptx", MOCK / "turn2" / "result.pptx")
    shutil.copy2(GOLD / "final.pptx", MOCK / "final" / "result.pptx")
    write_manifest()


if __name__ == "__main__":
    main()
