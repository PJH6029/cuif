from __future__ import annotations

import shutil

from PIL import Image, ImageDraw
import yaml
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from cuif_eval.runner import run_task
from cuif_eval.pptx.extract import extract_charts, extract_images, extract_shapes, extract_text_by_slide, slide_count


def _write_chart_deck(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ["Germany", "India", "United States", "China"]
    data.add_series("2015", [30.0, 14.9, 13.6, 23.9])
    data.add_series("2021", [39.8, 19.1, 20.3, 28.4])
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(8), Inches(4.5), data)
    chart_shape.name = "renewable_share_clustered_chart"
    prs.save(path)


def _write_image_formula_deck(path, image_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = slide.shapes.add_picture(str(image_path), Inches(1), Inches(1.2), width=Inches(4.0))
    pic.name = "paper_architecture_figure"
    box = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(8.0), Inches(0.5))
    box.name = "attention_formula_box"
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Attention(Q, K, V) = softmax(QK^T / sqrt(d_k))V"
    prs.save(path)


def _write_reference_image(path):
    image = Image.new("RGB", (360, 220), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((30, 30, 140, 92), outline="#1F4E79", width=5)
    draw.rectangle((210, 30, 330, 92), outline="#2A9D8F", width=5)
    draw.line((140, 61, 210, 61), fill="#172033", width=4)
    draw.text((54, 52), "Encoder", fill="#172033")
    draw.text((240, 52), "Decoder", fill="#172033")
    image.save(path)


def _write_template_shape_deck(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.8))
    band.name = "brand_header_band"
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    band.line.color.rgb = RGBColor(0xFF, 0xBF, 0x00)
    prs.save(path)


def test_pptx_extraction_reads_slide_text_bbox_and_style(toy_task):
    pptx = toy_task / "mock_outputs" / "final" / "result.pptx"
    assert slide_count(pptx) == 2
    text = extract_text_by_slide(pptx)
    assert any("CUIF PoC Dashboard" in " ".join(values) for values in text.values())
    title = [s for s in extract_shapes(pptx) if "CUIF PoC Dashboard" in s.text][0]
    assert 0.05 <= title.x <= 0.15
    assert any(run.font_color == "#1F4E79" and run.bold for run in title.runs)


def test_pptx_chart_data_check_reads_native_chart_series(tmp_path):
    task = tmp_path / "chart_task"
    (task / "artifacts").mkdir(parents=True)
    (task / "mock_outputs" / "turn1").mkdir(parents=True)
    _write_chart_deck(task / "artifacts" / "seed.pptx")
    shutil.copy2(task / "artifacts" / "seed.pptx", task / "mock_outputs" / "turn1" / "result.pptx")
    (task / "manifest.yaml").write_text(
        yaml.safe_dump(
            {
                "manifest_version": "0.1",
                "id": "chart_task",
                "title": "Chart data smoke task",
                "primary_artifact_family": "pptx",
                "artifact_families": ["pptx"],
                "tracks": ["open_tool"],
                "artifacts": {
                    "package": {"seed": {"path": "artifacts/seed.pptx", "type": "pptx", "role": "seed"}},
                    "expected_outputs": {"turn1": {"result": {"path": "result.pptx", "type": "pptx"}}},
                },
                "turns": [
                    {
                        "id": "turn1",
                        "instruction": "Create the chart.",
                        "expected_output": "result",
                        "checks": [
                            {"id": "exists", "evaluator": "file_exists", "artifact": "run.outputs.turn1.result", "points": 1},
                            {
                                "id": "chart_data",
                                "evaluator": "pptx_chart_data",
                                "artifact": "run.outputs.turn1.result",
                                "points": 4,
                                "params": {
                                    "selector": {"slide": 1, "name": "renewable_share_clustered_chart"},
                                    "chart_type": "COLUMN_CLUSTERED",
                                    "categories": ["Germany", "India", "United States", "China"],
                                    "series": [
                                        {"name": "2015", "values": [30.0, 14.9, 13.6, 23.9]},
                                        {"name": "2021", "values": [39.8, 19.1, 20.3, 28.4]},
                                    ],
                                    "value_tolerance": 0.01,
                                },
                            },
                        ],
                    }
                ],
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )

    charts = extract_charts(task / "mock_outputs" / "turn1" / "result.pptx")
    assert charts[0].categories == ["Germany", "India", "United States", "China"]

    result = run_task(task, adapter_name="mock", out=tmp_path / "run", skip_judges=True)
    statuses = {r.check_id: r.status for r in result["results"]}
    assert statuses["chart_data"] == "pass"


def test_pptx_image_and_formula_checks(tmp_path):
    task = tmp_path / "image_formula_task"
    (task / "artifacts" / "inputs").mkdir(parents=True)
    (task / "mock_outputs" / "turn1").mkdir(parents=True)
    reference = task / "artifacts" / "inputs" / "reference_figure.png"
    _write_reference_image(reference)
    _write_image_formula_deck(task / "artifacts" / "seed.pptx", reference)
    shutil.copy2(task / "artifacts" / "seed.pptx", task / "mock_outputs" / "turn1" / "result.pptx")
    (task / "manifest.yaml").write_text(
        yaml.safe_dump(
            {
                "manifest_version": "0.1",
                "id": "image_formula_task",
                "title": "Image and formula smoke task",
                "primary_artifact_family": "pptx",
                "artifact_families": ["pptx"],
                "tracks": ["open_tool"],
                "artifacts": {
                    "package": {
                        "seed": {"path": "artifacts/seed.pptx", "type": "pptx", "role": "seed"},
                        "reference_figure": {"path": "artifacts/inputs/reference_figure.png", "type": "png", "role": "source_input"},
                    },
                    "expected_outputs": {"turn1": {"result": {"path": "result.pptx", "type": "pptx"}}},
                },
                "turns": [
                    {
                        "id": "turn1",
                        "new_inputs": {"textual": [], "visual": ["package.reference_figure"]},
                        "instruction": "Embed the figure and formula.",
                        "expected_output": "result",
                        "checks": [
                            {"id": "exists", "evaluator": "file_exists", "artifact": "run.outputs.turn1.result", "points": 1},
                            {
                                "id": "figure_count",
                                "evaluator": "pptx_image_count",
                                "artifact": "run.outputs.turn1.result",
                                "points": 1,
                                "params": {"selector": {"slide": 1}, "min_count": 1},
                            },
                            {
                                "id": "figure_match",
                                "evaluator": "pptx_image_match",
                                "artifact": "run.outputs.turn1.result",
                                "points": 3,
                                "params": {
                                    "source": "package.reference_figure",
                                    "selector": {"slide": 1, "name": "paper_architecture_figure"},
                                    "min_similarity": 0.98,
                                    "region": {"slide": 1, "x_min": 0.05, "y_min": 0.10, "x_max": 0.56, "y_max": 0.52},
                                    "tolerance": 0.02,
                                },
                            },
                            {
                                "id": "formula_present",
                                "evaluator": "pptx_formula_present",
                                "artifact": "run.outputs.turn1.result",
                                "points": 2,
                                "params": {
                                    "formula": "Attention(Q,K,V)=softmax(QK^T/sqrt(d_k))V",
                                    "selector": {"slide": 1, "name": "attention_formula_box"},
                                    "region": {"slide": 1, "x_min": 0.05, "y_min": 0.65, "x_max": 0.92, "y_max": 0.88},
                                    "tolerance": 0.02,
                                },
                            },
                        ],
                    }
                ],
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )

    images = extract_images(task / "mock_outputs" / "turn1" / "result.pptx")
    assert images[0].shape_name == "paper_architecture_figure"

    result = run_task(task, adapter_name="mock", out=tmp_path / "run", skip_judges=True)
    statuses = {r.check_id: r.status for r in result["results"]}
    assert statuses["figure_count"] == "pass"
    assert statuses["figure_match"] == "pass"
    assert statuses["formula_present"] == "pass"


def test_pptx_style_check_scores_shape_fill_and_line_colors(tmp_path):
    task = tmp_path / "shape_style_task"
    (task / "artifacts").mkdir(parents=True)
    (task / "mock_outputs" / "turn1").mkdir(parents=True)
    _write_template_shape_deck(task / "artifacts" / "seed.pptx")
    shutil.copy2(task / "artifacts" / "seed.pptx", task / "mock_outputs" / "turn1" / "result.pptx")
    (task / "manifest.yaml").write_text(
        yaml.safe_dump(
            {
                "manifest_version": "0.1",
                "id": "shape_style_task",
                "title": "Shape style smoke task",
                "primary_artifact_family": "pptx",
                "artifact_families": ["pptx"],
                "tracks": ["open_tool"],
                "artifacts": {
                    "package": {"seed": {"path": "artifacts/seed.pptx", "type": "pptx", "role": "seed"}},
                    "expected_outputs": {"turn1": {"result": {"path": "result.pptx", "type": "pptx"}}},
                },
                "turns": [
                    {
                        "id": "turn1",
                        "instruction": "Keep the template header band style.",
                        "expected_output": "result",
                        "checks": [
                            {"id": "exists", "evaluator": "file_exists", "artifact": "run.outputs.turn1.result", "points": 1},
                            {
                                "id": "header_band_style",
                                "evaluator": "pptx_style_check",
                                "artifact": "run.outputs.turn1.result",
                                "points": 4,
                                "params": {
                                    "selector": {"slide": 1, "name": "brand_header_band"},
                                    "fill_color": "#1F4E79",
                                    "line_color": "#FFBF00",
                                },
                            },
                            {
                                "id": "header_band_wrong_fill",
                                "evaluator": "pptx_style_check",
                                "artifact": "run.outputs.turn1.result",
                                "points": 1,
                                "params": {
                                    "selector": {"slide": 1, "name": "brand_header_band"},
                                    "fill_color": "#FFFFFF",
                                },
                            },
                        ],
                    }
                ],
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )

    result = run_task(task, adapter_name="mock", out=tmp_path / "run", skip_judges=True)
    statuses = {r.check_id: r.status for r in result["results"]}
    assert statuses["header_band_style"] == "pass"
    assert statuses["header_band_wrong_fill"] == "fail"


def test_passing_toy_task_has_required_pptx_checks_pass(toy_task, tmp_path):
    result = run_task(toy_task, adapter_name="mock", out=tmp_path / "run", skip_judges=True)
    statuses = {r.check_id: r.status for r in result["results"]}
    assert statuses["title_present"] == "pass"
    assert statuses["title_region"] == "pass"
    assert statuses["brand_title_style"] == "pass"
    assert statuses["slide2_preserved"] == "pass"
    assert statuses["turn1_rendered_review"] in {"pass", "skipped", "blocked"}
    assert statuses["optional_llm_summary_judge"] == "skipped"


def test_broken_fixture_catches_text_style_layout_and_preservation_failures(toy_task, tmp_path):
    result = run_task(toy_task, adapter_name="mock", out=tmp_path / "broken", skip_judges=True, adapter_config={"mock_outputs_dir": "mock_outputs_broken"})
    statuses = {r.check_id: r.status for r in result["results"]}
    assert statuses["final_title_present"] == "fail"
    assert statuses["brand_title_style"] == "blocked"
    assert statuses["title_layout_preserved"] == "blocked"
    assert statuses["slide2_preserved"] == "fail"
