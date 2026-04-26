from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TASK = ROOT / "poc" / "tasks" / "toy_pptx_layout"


def add_textbox(slide, left, top, width, height, text, *, size=24, color="172033", bold=False, name=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if name:
        box.name = name
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color.replace("#", ""))
    return box


def add_panel(slide, left, top, width, height, title, body):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(204, 213, 224)
    add_textbox(slide, left + 0.2, top + 0.2, width - 0.4, 0.35, title, size=18, color="1F4E79", bold=True)
    add_textbox(slide, left + 0.2, top + 0.75, width - 0.4, 0.7, body, size=14, color="172033")


def base_deck(title="Seed CUIF PoC", subtitle="Prepared seed deck", *, brand=False, damage_slide2=False):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank)
    bg = slide1.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(247, 251, 255)
    add_textbox(slide1, 0.7, 0.35, 8.6, 0.55, title, size=34 if brand else 30, color="1F4E79" if brand else "172033", bold=brand, name="dashboard-title")
    add_textbox(slide1, 0.75, 1.0, 8.5, 0.35, subtitle, size=18, color="44546A")
    add_panel(slide1, 0.8, 1.65, 3.8, 2.0, "Task loop", "Manifest → adapter → outputs → evaluator → report")
    add_panel(slide1, 5.3, 1.65, 3.8, 2.0, "Checks", "Text, layout, style, preservation, optional judges")

    slide2 = prs.slides.add_slide(blank)
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    add_textbox(slide2, 0.7, 0.45, 8.7, 0.45, "Protected context slide", size=26, color="1F4E79", bold=True)
    invariant = "Do not edit: CUIF benchmark invariant"
    if damage_slide2:
        invariant = "Damaged protected context"
    add_textbox(slide2, 0.7, 1.25, 8.7, 0.45, invariant, size=20, color="172033", name="protected-invariant")
    add_textbox(slide2, 0.7, 2.0, 8.7, 0.45, "This slide simulates collateral-damage preservation requirements.", size=16, color="44546A")
    return prs


def save(prs, rel):
    path = TASK / rel
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    print(path)


def main():
    save(base_deck(), "artifacts/seed.pptx")
    turn1 = base_deck("CUIF PoC Dashboard", "Partial-credit office evaluation")
    save(turn1, "artifacts/gold/turn1.pptx")
    save(turn1, "mock_outputs/turn1/result.pptx")
    final = base_deck("CUIF PoC Dashboard", "Partial-credit office evaluation", brand=True)
    save(final, "artifacts/gold/final.pptx")
    save(final, "mock_outputs/final/result.pptx")
    broken = base_deck("Wrong Dashboard", "Partial-credit office evaluation", brand=False, damage_slide2=True)
    save(turn1, "mock_outputs_broken/turn1/result.pptx")
    save(broken, "mock_outputs_broken/final/result.pptx")


if __name__ == "__main__":
    main()
